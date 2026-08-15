from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt


SOURCE = Path("/home/xckj/suyuan/backend/backend_data_registry/uploads/b233f437-c359-4c18-b777-33128d8b25ae.pptx")
OUT_DIR = Path("/home/xckj/suyuan/backend/backend_data_registry/pptx_templates/air_quality_formal_v1")
OUTPUT = OUT_DIR / "air_quality_formal_v1.pptx"
PLACEHOLDER = OUT_DIR / "assets" / "visual_placeholder.png"

BLUE = RGBColor(31, 88, 190)
LIGHT_BLUE = RGBColor(231, 240, 255)
MID_BLUE = RGBColor(86, 145, 220)
TEXT = RGBColor(36, 48, 64)
MUTED = RGBColor(102, 116, 139)
LINE = RGBColor(180, 205, 238)
WHITE = RGBColor(255, 255, 255)


SLIDES = [
    {
        "kind": "cover",
        "title": "{{cover.title}}",
        "subtitle": "{{cover.subtitle}}",
        "footer": "{{cover.organization}} | {{cover.date}}",
    },
    {
        "kind": "toc",
        "title": "目录",
        "items": ["{{toc.item_1}}", "{{toc.item_2}}", "{{toc.item_3}}", "{{toc.item_4}}"],
    },
    {
        "kind": "summary",
        "title": "{{executive_summary.title}}",
        "message": "{{executive_summary.message}}",
        "items": [
            ("核心判断", "{{executive_summary.finding_1}}"),
            ("关键变化", "{{executive_summary.finding_2}}"),
            ("管理建议", "{{executive_summary.finding_3}}"),
        ],
    },
    {
        "kind": "process",
        "title": "{{framework.title}}",
        "steps": ["{{framework.step_1}}", "{{framework.step_2}}", "{{framework.step_3}}", "{{framework.step_4}}"],
    },
    {
        "kind": "visual_full",
        "title": "{{architecture.title}}",
        "label": "{{architecture.caption}}",
    },
    {
        "kind": "metric_dashboard",
        "title": "{{metric_dashboard.title}}",
        "metrics": [
            ("{{metric_1.label}}", "{{metric_1.value}}", "{{metric_1.note}}"),
            ("{{metric_2.label}}", "{{metric_2.value}}", "{{metric_2.note}}"),
            ("{{metric_3.label}}", "{{metric_3.value}}", "{{metric_3.note}}"),
            ("{{metric_4.label}}", "{{metric_4.value}}", "{{metric_4.note}}"),
        ],
    },
    {
        "kind": "map_insight",
        "title": "{{map_insight.title}}",
        "items": ["{{map_insight.finding_1}}", "{{map_insight.finding_2}}", "{{map_insight.finding_3}}"],
    },
    {
        "kind": "chart_insight",
        "title": "{{chart_insight.title}}",
        "items": ["{{chart_insight.finding_1}}", "{{chart_insight.finding_2}}", "{{chart_insight.finding_3}}"],
    },
    {
        "kind": "ranking_table",
        "title": "{{ranking.title}}",
        "headers": ["{{table.header_1}}", "{{table.header_2}}", "{{table.header_3}}", "{{table.header_4}}"],
    },
    {
        "kind": "evidence_dashboard",
        "title": "{{evidence.title}}",
        "captions": ["{{evidence.caption_1}}", "{{evidence.caption_2}}", "{{evidence.caption_3}}"],
    },
    {
        "kind": "closed_loop",
        "title": "{{closed_loop.title}}",
        "steps": ["{{closed_loop.step_1}}", "{{closed_loop.step_2}}", "{{closed_loop.step_3}}", "{{closed_loop.step_4}}"],
    },
    {
        "kind": "actions",
        "title": "{{actions.title}}",
        "items": [
            ("近期措施", "{{actions.near_term}}"),
            ("重点区域", "{{actions.region_focus}}"),
            ("跟踪评估", "{{actions.follow_up}}"),
        ],
    },
    {
        "kind": "comparison",
        "title": "{{comparison.title}}",
        "left": "{{comparison.left_title}}",
        "right": "{{comparison.right_title}}",
    },
    {
        "kind": "report_assistant",
        "title": "{{report.title}}",
        "items": ["{{report.point_1}}", "{{report.point_2}}", "{{report.point_3}}"],
    },
    {
        "kind": "thanks",
        "title": "谢谢观看！",
        "subtitle": "{{thanks.subtitle}}",
        "footer": "{{thanks.contact}}",
    },
]


def emu_to_in(value: int) -> float:
    return value / 914400


def make_placeholder_image() -> None:
    PLACEHOLDER.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", (1600, 900), (236, 244, 255))
    draw = ImageDraw.Draw(img)
    for i in range(0, 1600, 80):
        draw.line((i, 0, i - 500, 900), fill=(215, 229, 250), width=3)
    draw.rectangle((20, 20, 1580, 880), outline=(88, 145, 220), width=8)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 52)
    except Exception:
        font = ImageFont.load_default()
    text = "VISUAL / CHART / MAP"
    bbox = draw.textbbox((0, 0), text, font=font)
    draw.text(((1600 - (bbox[2] - bbox[0])) / 2, 410), text, fill=(50, 90, 150), font=font)
    img.save(PLACEHOLDER)


def remove_shape(shape) -> None:
    element = shape._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def is_logo(shape) -> bool:
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    x = emu_to_in(shape.left)
    y = emu_to_in(shape.top)
    w = emu_to_in(shape.width)
    h = emu_to_in(shape.height)
    return y < 0.75 and w <= 1.25 and h <= 0.55 and (x < 1.2 or x > 11.2)


def clean_slide(slide, keep_large_images: bool = False) -> None:
    for shape in list(slide.shapes):
        if is_logo(shape):
            continue
        if keep_large_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        y = emu_to_in(getattr(shape, "top", 0) or 0)
        h = emu_to_in(getattr(shape, "height", 0) or 0)
        if y < 0.78 and h < 0.35 and not getattr(shape, "has_text_frame", False):
            continue
        remove_shape(shape)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 16,
             color: RGBColor = TEXT, bold: bool = False, align: str | None = None,
             fill: RGBColor | None = None, line: RGBColor | None = None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill or line:
        box.fill.solid()
        box.fill.fore_color.rgb = fill or WHITE
        box.line.color.rgb = line or fill or WHITE
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = 2
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str) -> None:
    add_text(slide, title, 0.78, 0.18, 9.4, 0.45, 19, TEXT, True)


def add_visual(slide, x: float, y: float, w: float, h: float):
    return slide.shapes.add_picture(str(PLACEHOLDER), Inches(x), Inches(y), Inches(w), Inches(h))


def add_card(slide, title: str, body: str, x: float, y: float, w: float, h: float) -> None:
    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(248, 251, 255)
    rect.line.color.rgb = LINE
    add_text(slide, title, x + 0.12, y + 0.12, w - 0.24, 0.32, 13, BLUE, True)
    add_text(slide, body, x + 0.12, y + 0.54, w - 0.24, h - 0.66, 11, MUTED)


def add_metric(slide, label: str, value: str, note: str, x: float, y: float, w: float, h: float) -> None:
    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(246, 250, 255)
    rect.line.color.rgb = LINE
    add_text(slide, label, x + 0.15, y + 0.12, w - 0.3, 0.28, 11, MUTED, True, "center")
    add_text(slide, value, x + 0.15, y + 0.52, w - 0.3, 0.55, 24, BLUE, True, "center")
    add_text(slide, note, x + 0.15, y + 1.13, w - 0.3, 0.3, 10, MUTED, False, "center")


def add_template_content(slide, spec: dict) -> None:
    kind = spec["kind"]
    if kind == "cover":
        clean_slide(slide, keep_large_images=True)
    else:
        clean_slide(slide)
        add_title(slide, spec["title"])

    if kind == "cover":
        add_text(slide, spec["title"], 1.1, 2.35, 6.3, 0.62, 24, WHITE, True, "center", BLUE, BLUE)
        add_text(slide, spec["subtitle"], 1.1, 3.02, 6.3, 0.42, 15, WHITE, False, "center", BLUE, BLUE)
        add_text(slide, spec["footer"], 0.62, 6.45, 4.2, 0.28, 9, MUTED)
    elif kind == "toc":
        for idx, item in enumerate(spec["items"], start=1):
            y = 1.55 + (idx - 1) * 1.02
            add_text(slide, f"{idx:02d}", 4.45, y, 0.72, 0.36, 13, WHITE, True, "center", BLUE, BLUE)
            add_text(slide, item, 5.45, y - 0.02, 4.8, 0.42, 14, BLUE, True)
    elif kind == "summary":
        add_text(slide, spec["message"], 0.9, 1.25, 11.4, 0.6, 17, BLUE, True, "center", LIGHT_BLUE, LINE)
        for i, item in enumerate(spec["items"]):
            add_card(slide, item[0], item[1], 0.9 + i * 4.05, 2.2, 3.65, 2.4)
        add_visual(slide, 0.9, 5.0, 11.4, 1.35)
    elif kind == "process":
        for i, step in enumerate(spec["steps"]):
            x = 0.95 + i * 3.05
            add_text(slide, str(i + 1), x, 2.0, 0.48, 0.42, 15, WHITE, True, "center", BLUE, BLUE)
            add_card(slide, f"步骤 {i + 1}", step, x, 2.6, 2.45, 2.1)
        add_visual(slide, 0.95, 5.35, 11.45, 0.85)
    elif kind == "visual_full":
        add_visual(slide, 0.95, 1.25, 11.45, 4.95)
        add_text(slide, spec["label"], 0.95, 6.35, 11.45, 0.28, 10, MUTED)
    elif kind == "metric_dashboard":
        for i, metric in enumerate(spec["metrics"]):
            add_metric(slide, *metric, 0.95 + i * 3.05, 1.58, 2.45, 1.75)
        add_visual(slide, 0.95, 3.7, 11.45, 2.55)
    elif kind in {"map_insight", "chart_insight"}:
        add_visual(slide, 0.95, 1.35, 7.0, 4.95)
        for i, item in enumerate(spec["items"], start=1):
            add_card(slide, f"洞察 {i}", item, 8.35, 1.35 + (i - 1) * 1.65, 3.8, 1.25)
    elif kind == "ranking_table":
        rows = 5
        cols = 4
        x0, y0, w, h = 0.95, 1.45, 11.45, 4.55
        table = slide.shapes.add_table(rows, cols, Inches(x0), Inches(y0), Inches(w), Inches(h)).table
        for c in range(cols):
            table.cell(0, c).text = spec["headers"][c]
        for r in range(1, rows):
            for c in range(cols):
                table.cell(r, c).text = f"{{{{table.r{r}c{c + 1}}}}}"
    elif kind == "evidence_dashboard":
        boxes = [(0.95, 1.35), (5.05, 1.35), (9.15, 1.35)]
        for i, (x, y) in enumerate(boxes):
            add_visual(slide, x, y, 3.25, 3.15)
            add_text(slide, spec["captions"][i], x, y + 3.28, 3.25, 0.35, 10, MUTED, False, "center")
        add_text(slide, "{{evidence.key_message}}", 0.95, 5.45, 11.45, 0.58, 15, BLUE, True, "center", LIGHT_BLUE, LINE)
    elif kind == "closed_loop":
        for i, step in enumerate(spec["steps"]):
            x = 1.0 + i * 3.0
            add_text(slide, f"{i + 1}", x, 2.2, 0.5, 0.5, 17, WHITE, True, "center", BLUE, BLUE)
            add_card(slide, f"环节 {i + 1}", step, x, 3.0, 2.4, 1.8)
        add_visual(slide, 1.0, 5.25, 11.0, 0.9)
    elif kind == "actions":
        for i, item in enumerate(spec["items"]):
            add_card(slide, item[0], item[1], 1.0 + i * 4.0, 1.55, 3.45, 3.9)
        add_visual(slide, 1.0, 5.75, 11.45, 0.45)
    elif kind == "comparison":
        add_card(slide, spec["left"], "{{comparison.left_body}}", 0.95, 1.45, 5.45, 4.8)
        add_card(slide, spec["right"], "{{comparison.right_body}}", 6.95, 1.45, 5.45, 4.8)
        add_visual(slide, 5.92, 2.15, 0.55, 3.1)
    elif kind == "report_assistant":
        add_visual(slide, 0.95, 1.35, 6.6, 4.65)
        for i, item in enumerate(spec["items"], start=1):
            add_card(slide, f"要点 {i}", item, 8.0, 1.35 + (i - 1) * 1.55, 3.85, 1.16)
    elif kind == "thanks":
        add_visual(slide, 7.25, 1.35, 4.65, 4.55)
        add_text(slide, spec["title"], 1.25, 2.75, 5.2, 0.62, 23, WHITE, True, "center", BLUE, BLUE)
        add_text(slide, spec["subtitle"], 1.25, 3.38, 5.2, 0.36, 13, WHITE, False, "center", BLUE, BLUE)
        add_text(slide, spec["footer"], 0.62, 6.55, 5.0, 0.28, 8, MUTED)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    make_placeholder_image()
    prs = Presentation(str(SOURCE))
    for slide, spec in zip(prs.slides, SLIDES):
        add_template_content(slide, spec)
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
