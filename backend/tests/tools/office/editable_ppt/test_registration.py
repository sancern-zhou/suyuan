from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from app.agent.prompts.tool_registry import ASSISTANT_TOOL_NAMES
from app.agent.skill_metadata import SKILL_METADATA
from app.tools import create_global_tool_registry
from app.tools.office.ppt_master_tool import CreatePptxWithPptMasterTool


def _run_typeface(run, element_name):
    run_properties = run._r.get_or_add_rPr()
    font_element = run_properties.find(qn(element_name))
    return None if font_element is None else font_element.get("typeface")


def test_editable_ppt_tool_is_registered_and_exposed():
    registry = create_global_tool_registry()
    assert registry.get_tool("manage_editable_ppt") is not None
    assert "manage_editable_ppt" in ASSISTANT_TOOL_NAMES
    assert "create_pptx_with_ppt_master" in ASSISTANT_TOOL_NAMES


def test_workflow_declares_direct_file_edit_and_strict_gate():
    guide = Path("app/tools/office/editable_ppt/references/workflow.md").read_text(encoding="utf-8")
    assert "edit_file" in guide
    assert "3–5" in guide
    assert "strict" in guide
    assert "无需重新生成" in guide


def test_workflow_contains_copyable_native_element_contracts():
    guide = Path("app/tools/office/editable_ppt/references/workflow.md").read_text(encoding="utf-8")
    assert 'data-pptx-ref="roi-chart"' in guide
    assert 'kind: "chart"' in guide
    assert 'categories: ["Q1", "Q2"]' in guide
    assert 'series: [{ name: "节省工时", values: [120, 260] }]' in guide
    assert 'kind: "table"' in guide
    assert 'data: { rows:' in guide
    assert "示例数据" in guide


def test_editable_ppt_generation_scenario_metadata_is_selectable():
    metadata = SKILL_METADATA["editable_ppt_generation"]
    assert {"manage_editable_ppt", "read_file", "edit_file"} <= set(metadata["required_tools"])


def test_ppt_master_text_sets_latin_and_east_asian_typefaces():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    text_box = CreatePptxWithPptMasterTool()._add_text(
        slide,
        "中文标题",
        1,
        1,
        4,
        1,
        font_size=20,
        color="000000",
    )

    run = text_box.text_frame.paragraphs[0].runs[0]
    assert _run_typeface(run, "a:latin") == "Microsoft YaHei"
    assert _run_typeface(run, "a:ea") == "Microsoft YaHei"


def test_ppt_master_table_sets_latin_and_east_asian_typefaces():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        1,
        1,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1),
    ).table

    CreatePptxWithPptMasterTool()._format_table_cell(
        table.cell(0, 0),
        "中文单元格",
        font_size=12,
        color="000000",
        fill="FFFFFF",
        bold=False,
        RGBColor=RGBColor,
        Pt=Pt,
    )

    run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    assert _run_typeface(run, "a:latin") == "Microsoft YaHei"
    assert _run_typeface(run, "a:ea") == "Microsoft YaHei"
