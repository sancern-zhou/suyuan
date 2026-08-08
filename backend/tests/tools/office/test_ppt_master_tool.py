from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.tools.office.ppt_master_tool import CreatePptxWithPptMasterTool


@pytest.mark.asyncio
async def test_ppt_master_tool_creates_project_and_diverse_layouts(tmp_path: Path):
    result = await CreatePptxWithPptMasterTool().execute(
        title="广东省空气质量会商分析",
        purpose="government_briefing",
        audience="生态环境管理部门",
        style="government_consulting",
        outline=[
            {"title": "核心结论", "points": ["PM2.5持续改善", "臭氧仍需重点关注"]},
            {"title": "趋势变化", "chart": {"type": "line"}, "points": ["1-5月同比改善"]},
            {"title": "城市排名", "chart": {"type": "bar"}, "points": ["重点城市差异明显"]},
            {"title": "下一步建议", "points": ["强化区域联防联控", "推进重点源精细管控"]},
        ],
        output_file=str(tmp_path / "air_quality.pptx"),
        project_dir=str(tmp_path / "project"),
        enable_preview=False,
        run_validation=False,
    )

    assert result["success"] is True
    assert Path(result["data"]["file_path"]).exists()
    assert Path(result["data"]["project_dir"]).exists()
    assert Path(result["data"]["design_spec_path"]).exists()
    assert Path(result["data"]["spec_lock_path"]).exists()
    assert len(result["data"]["page_plan"]) == 5

    layouts = [page["layout"] for page in result["data"]["page_plan"]]
    assert len(set(layouts[1:])) >= 3
    assert "image_left_text_right" not in layouts
    assert result["data"]["workflow"] == "ppt_master"
    assert result["data"]["quality_gate"]["status"] in {"pass", "warning", "rewrite_required"}
    assert result["refs"]["files"][0]["path"] == result["data"]["file_path"]
    assert result["refs"]["artifacts"][0]["file_path"] == result["data"]["file_path"]
    assert result["llm_resume"]["artifact_path"] == result["data"]["file_path"]
    assert result["llm_resume"]["project_dir"] == result["data"]["project_dir"]
    assert "present_artifact" in result["llm_resume"]["tool_hint"]


def test_ppt_master_tool_schema_describes_production_workflow():
    schema = CreatePptxWithPptMasterTool().get_function_schema()

    assert schema["name"] == "create_pptx_with_ppt_master"
    properties = schema["parameters"]["properties"]
    assert {"title", "purpose", "outline"}.issubset(properties)
    assert "DeckSpec" not in schema["description"]
    outline_description = properties["outline"]["description"]
    assert "chart.image_path" in outline_description
    assert "生成时直接插入" in outline_description


@pytest.mark.asyncio
async def test_ppt_master_tool_renders_native_table_from_slide_plan(tmp_path: Path):
    result = await CreatePptxWithPptMasterTool().execute(
        title="表格绘制测试",
        slide_plan=[
            {
                "title": "关键指标对比",
                "message": "表格应作为 PowerPoint 原生表格生成",
                "shapes": [
                    {
                        "type": "table",
                        "x": 0.8,
                        "y": 1.4,
                        "w": 11.6,
                        "h": 2.1,
                        "rows": [
                            ["城市", "PM2.5", "同比"],
                            ["广州", "28", "-6%"],
                            ["深圳", "22", "-8%"],
                        ],
                        "font_size": 12,
                        "header_fill": "174A7C",
                        "header_color": "FFFFFF",
                    }
                ],
            }
        ],
        output_file=str(tmp_path / "native_table.pptx"),
        project_dir=str(tmp_path / "project_table"),
        enable_preview=False,
        run_validation=False,
    )

    assert result["success"] is True
    prs = Presentation(result["data"]["file_path"])
    content_slide = prs.slides[1]
    tables = [shape.table for shape in content_slide.shapes if getattr(shape, "has_table", False)]

    assert len(tables) == 1
    table = tables[0]
    assert len(table.rows) == 3
    assert len(table.columns) == 3
    assert table.cell(0, 0).text == "城市"
    assert table.cell(1, 2).text == "-6%"


@pytest.mark.asyncio
async def test_ppt_master_tool_inserts_real_chart_image_when_provided(tmp_path: Path):
    chart_image = tmp_path / "actual_chart.png"
    Image.new("RGB", (640, 360), "#2f6fed").save(chart_image)

    result = await CreatePptxWithPptMasterTool().execute(
        title="图表插入测试",
        outline=[
            {
                "title": "真实图表页",
                "chart": {"type": "image", "image_path": str(chart_image)},
                "points": ["应插入实际图表图片"],
            }
        ],
        output_file=str(tmp_path / "chart_image.pptx"),
        project_dir=str(tmp_path / "project_image"),
        enable_preview=False,
        run_validation=False,
    )

    assert result["success"] is True
    prs = Presentation(result["data"]["file_path"])
    content_slide = prs.slides[1]
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in content_slide.shapes)
    assert result["data"]["page_plan"][1]["chart"]["resolved_asset_path"] == str(chart_image)


@pytest.mark.asyncio
async def test_ppt_master_tool_renders_native_chart_from_chart_data(tmp_path: Path):
    result = await CreatePptxWithPptMasterTool().execute(
        title="原生图表测试",
        outline=[
            {
                "title": "真实数据图表页",
                "chart": {
                    "type": "bar",
                    "categories": ["广州", "深圳", "佛山"],
                    "series": [{"name": "AQI", "values": [52, 48, 61]}],
                },
                "points": ["应生成 PowerPoint 原生图表"],
            }
        ],
        output_file=str(tmp_path / "native_chart.pptx"),
        project_dir=str(tmp_path / "project_native"),
        enable_preview=False,
        run_validation=False,
    )

    assert result["success"] is True
    prs = Presentation(result["data"]["file_path"])
    content_slide = prs.slides[1]
    assert any(shape.shape_type == MSO_SHAPE_TYPE.CHART for shape in content_slide.shapes)
    assert result["data"]["page_plan"][1]["chart"]["render_mode"] == "native_chart"


@pytest.mark.asyncio
async def test_ppt_master_tool_deduplicates_cover_and_uses_agenda_layout(tmp_path: Path):
    result = await CreatePptxWithPptMasterTool().execute(
        title="年度汇报",
        outline=[
            {"title": "封面", "points": ["年度汇报"]},
            {"title": "目录", "points": ["背景", "分析", "建议"]},
            {"title": "背景", "points": ["业务背景"]},
        ],
        output_file=str(tmp_path / "agenda.pptx"),
        project_dir=str(tmp_path / "project_agenda"),
        enable_preview=False,
        run_validation=False,
    )

    assert result["success"] is True
    page_plan = result["data"]["page_plan"]
    assert [page["title"] for page in page_plan].count("封面") == 0
    assert page_plan[0]["role"] == "cover"
    assert page_plan[1]["role"] == "agenda"
    assert page_plan[1]["layout"] == "agenda"


@pytest.mark.asyncio
async def test_ppt_master_tool_varies_repeated_chart_layouts(tmp_path: Path):
    outline = [
        {
            "title": f"图表页 {index}",
            "chart": {
                "type": "bar",
                "categories": ["A", "B", "C"],
                "series": [{"name": "值", "values": [index, index + 1, index + 2]}],
            },
            "points": [f"第 {index} 页洞察"],
        }
        for index in range(1, 7)
    ]
    result = await CreatePptxWithPptMasterTool().execute(
        title="图表版式测试",
        outline=outline,
        output_file=str(tmp_path / "chart_layouts.pptx"),
        project_dir=str(tmp_path / "project_chart_layouts"),
        enable_preview=False,
        run_validation=False,
    )

    chart_layouts = [page["layout"] for page in result["data"]["page_plan"] if page.get("chart")]
    assert len(set(chart_layouts)) >= 4
    assert max(
        sum(1 for item in group)
        for _, group in __import__("itertools").groupby(chart_layouts)
    ) <= 1


def test_ppt_master_quality_gate_includes_validation_failure():
    tool = CreatePptxWithPptMasterTool()
    page_plan = [
        {"slide": 1, "layout": "cover_statement", "role": "cover"},
        {"slide": 2, "layout": "card_grid", "role": "content"},
    ]
    gate = tool._workflow_quality_gate(
        page_plan,
        validation={
            "success": False,
            "overflow_issues": [{"slide": 1, "type": "rendered_content_overflow"}],
            "design_quality": {"issues": [{"slide": 2, "type": "text_only_slide"}]},
        },
    )

    assert gate["status"] == "rewrite_required"
    assert gate["rewrite_required"] is True
    assert any(issue["type"] == "validation_failed" for issue in gate["issues"])
    assert any(issue["type"] == "rendered_content_overflow" for issue in gate["issues"])


def test_template_generation_modules_are_removed():
    office_dir = Path(__file__).resolve().parents[3] / "app" / "tools" / "office"

    assert not (office_dir / "analyze_pptx_template_tool.py").exists()
    assert not (office_dir / "create_pptx_from_template_tool.py").exists()
