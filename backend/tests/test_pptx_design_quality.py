from pathlib import Path
import json

import pytest
from PIL import Image
from PIL import ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.agent.tool_adapter import call_llm_tool
from app.tools.office.ppt_master_tool import CreatePptxWithPptMasterTool
from app.tools.office.validate_pptx_tool import ValidatePptxTool


@pytest.mark.asyncio
async def test_ppt_master_tool_returns_structured_error_when_title_missing():
    result = await CreatePptxWithPptMasterTool().execute()

    assert result["success"] is False
    assert result["data"]["error"] == "title_required"
    assert "title 参数缺失" in result["summary"]


@pytest.mark.asyncio
async def test_ppt_master_llm_tool_call_handles_empty_args_without_type_error():
    result = await call_llm_tool("create_pptx_with_ppt_master")

    assert result["status"] == "failed"
    assert result["success"] is False
    assert result["data"]["error"] == "title_required"
    assert "missing 1 required positional argument" not in result.get("error", "")


def test_ppt_master_theme_normalizes_color_aliases():
    palette = CreatePptxWithPptMasterTool()._palette(
        "business_clean",
        {
            "primary_color": "#1E88E5",
            "secondary_color": "#43A047",
            "accent_color": "#FFA726",
        },
    )

    assert palette["primary"] == "1E88E5"
    assert palette["secondary"] == "43A047"
    assert palette["accent"] == "FFA726"


@pytest.mark.asyncio
async def test_ppt_master_tool_renders_agent_shape_plan_with_contained_image(tmp_path: Path):
    image_path = tmp_path / "wide_chart.png"
    Image.new("RGB", (800, 200), "#2f6fed").save(image_path)

    result = await CreatePptxWithPptMasterTool().execute(
        title="Agent绘制计划测试",
        slide_plan=[
            {
                "title": "动态图表页",
                "message": "Agent 自行决定图表和文本区域",
                "shapes": [
                    {
                        "type": "text",
                        "text": "动态图表页",
                        "x": 0.06,
                        "y": 0.06,
                        "w": 0.88,
                        "h": 0.08,
                        "unit": "relative",
                        "font_size": 26,
                        "bold": True,
                    },
                    {
                        "type": "image",
                        "path": str(image_path),
                        "x": 0.08,
                        "y": 0.22,
                        "w": 0.64,
                        "h": 0.42,
                        "unit": "relative",
                        "fit": "contain",
                    },
                    {
                        "type": "textbox",
                        "text": "图表保持比例，右侧洞察由 Agent 控制坐标。",
                        "x": 0.76,
                        "y": 0.24,
                        "w": 0.18,
                        "h": 0.24,
                        "unit": "relative",
                        "font_size": 14,
                    },
                ],
            }
        ],
        output_file=str(tmp_path / "agent_plan.pptx"),
        project_dir=str(tmp_path / "project_agent_plan"),
        enable_preview=False,
        run_validation=False,
    )

    assert result["success"] is True
    assert result["data"]["page_plan"][1]["layout"] == "agent_shape_plan"
    prs = Presentation(result["data"]["file_path"])
    slide = prs.slides[1]
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    assert pictures[0].width > pictures[0].height
    assert any(getattr(shape, "has_text_frame", False) and "动态图表页" in shape.text for shape in slide.shapes)


@pytest.mark.asyncio
async def test_ppt_master_shape_plan_cover_crops_image_inside_box(tmp_path: Path):
    image_path = tmp_path / "wide_chart.png"
    Image.new("RGB", (800, 200), "#2f6fed").save(image_path)

    result = await CreatePptxWithPptMasterTool().execute(
        title="Cover裁剪测试",
        slide_plan=[
            {
                "title": "封面图",
                "shapes": [
                    {
                        "type": "image",
                        "path": str(image_path),
                        "x": 0.1,
                        "y": 0.2,
                        "w": 0.3,
                        "h": 0.4,
                        "unit": "relative",
                        "fit": "cover",
                    }
                ],
            }
        ],
        output_file=str(tmp_path / "cover_crop.pptx"),
        project_dir=str(tmp_path / "project_cover_crop"),
        enable_preview=False,
        run_validation=False,
    )

    assert result["success"] is True
    prs = Presentation(result["data"]["file_path"])
    picture = next(shape for shape in prs.slides[1].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    expected_width = int(round(0.3 * 13.333 * 914400))
    expected_height = int(round(0.4 * 7.5 * 914400))
    assert abs(picture.width - expected_width) < 4
    assert abs(picture.height - expected_height) < 4
    assert picture.crop_left > 0
    assert picture.crop_right > 0


@pytest.mark.asyncio
async def test_ppt_master_applies_plan_patch_without_rewriting_untouched_slides(tmp_path: Path):
    base_plan = [
        {
            "slide": 1,
            "layout": "cover_statement",
            "role": "cover",
            "title": "基线PPT",
            "message": "面向决策的结构化汇报",
            "points": [],
        },
        {
            "slide": 2,
            "layout": "agent_shape_plan",
            "role": "content",
            "title": "保留页面",
            "message": "这页不应该被 Agent 重写",
            "points": [],
            "shapes": [
                {
                    "type": "text",
                    "text": "原始保留内容",
                    "x": 0.1,
                    "y": 0.2,
                    "w": 0.8,
                    "h": 0.1,
                    "unit": "relative",
                }
            ],
        },
        {
            "slide": 3,
            "layout": "agent_shape_plan",
            "role": "content",
            "title": "待拆分页",
            "message": "文字太密",
            "points": [],
            "shapes": [],
        },
    ]
    base_plan_path = tmp_path / "slide_plan.v1.json"
    base_plan_path.write_text(json.dumps(base_plan, ensure_ascii=False), encoding="utf-8")

    result = await CreatePptxWithPptMasterTool().execute(
        base_plan_path=str(base_plan_path),
        plan_patch={
            "replace_slides": [
                {
                    "slide": 3,
                    "slides": [
                        {
                            "title": "拆分后一",
                            "message": "保留原结论，降低密度",
                            "shapes": [
                                {
                                    "type": "text",
                                    "text": "拆分后一",
                                    "x": 0.1,
                                    "y": 0.2,
                                    "w": 0.8,
                                    "h": 0.1,
                                    "unit": "relative",
                                }
                            ],
                        },
                        {
                            "title": "拆分后二",
                            "message": "新增承接页",
                            "shapes": [
                                {
                                    "type": "text",
                                    "text": "拆分后二",
                                    "x": 0.1,
                                    "y": 0.2,
                                    "w": 0.8,
                                    "h": 0.1,
                                    "unit": "relative",
                                }
                            ],
                        },
                    ],
                }
            ]
        },
        output_file=str(tmp_path / "patched.pptx"),
        project_dir=str(tmp_path / "project_patched"),
        enable_preview=False,
        run_validation=False,
    )

    assert result["success"] is True
    data = result["data"]
    assert data["revision"]["base_plan_path"] == str(base_plan_path.resolve())
    assert data["slide_plan_path"].endswith("slide_plan.v2.json")
    revised_plan = json.loads(Path(data["slide_plan_path"]).read_text(encoding="utf-8"))
    assert [page["slide"] for page in revised_plan] == [1, 2, 3, 4]
    assert revised_plan[1]["title"] == "保留页面"
    assert revised_plan[1]["shapes"][0]["text"] == "原始保留内容"
    assert revised_plan[2]["title"] == "拆分后一"
    assert revised_plan[3]["title"] == "拆分后二"


def test_validate_pptx_design_quality_report(tmp_path: Path):
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(914400, 914400, 7315200, 914400)
    text_box.text = "设计质量测试"
    body = slide.shapes.add_textbox(914400, 1828800, 9144000, 3657600)
    body.text = "\n".join([f"这是一段用于模拟真实高文字密度的较长内容 {idx}" for idx in range(40)])
    pptx_path = tmp_path / "quality.pptx"
    presentation.save(pptx_path)

    report = ValidatePptxTool()._inspect_design_quality(pptx_path)

    assert report["enabled"] is True
    assert "score" in report
    assert report["grade"] in {"excellent", "good", "acceptable", "needs_improvement"}
    assert report["slides"][0]["text_lines"] >= 40
    assert any(issue["type"] == "high_text_density" for issue in report["issues"])


def test_validate_pptx_does_not_block_compact_kpi_cards_on_text_box_count_alone(tmp_path: Path):
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for index in range(15):
        box = slide.shapes.add_textbox(300000 + index * 30000, 300000 + index * 30000, 1200000, 300000)
        box.text = f"指标{index}：达标"
    pptx_path = tmp_path / "compact-kpis.pptx"
    presentation.save(pptx_path)

    report = ValidatePptxTool()._inspect_design_quality(pptx_path)

    assert report["slides"][0]["text_lines"] == 15
    assert not any(issue["type"] == "high_text_density" for issue in report["issues"])


def test_validate_pptx_counts_native_autoshapes_as_visual_elements(tmp_path: Path):
    pptx = pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 200000, 200000, 1200000, 800000)
    for index in range(3):
        box = slide.shapes.add_textbox(1800000, 300000 + index * 500000, 3000000, 400000)
        box.text = f"视觉卡片说明 {index}"
    pptx_path = tmp_path / "native-shapes.pptx"
    presentation.save(pptx_path)

    report = ValidatePptxTool()._inspect_design_quality(pptx_path)

    assert report["slides"][0]["visual_shapes"] == 1
    assert not any(issue["type"] == "text_only_slide" for issue in report["issues"])


def test_validate_pptx_builds_structured_factual_issues():
    tool = ValidatePptxTool()
    report = {
        "pages": [{"slide": 2, "png_path": "/tmp/page-002.png"}],
        "geometry": {"issues": [{"type": "shape_out_of_bounds", "slide": 2, "shape": 3}]},
        "rendered_overflow": {
            "issues": [
                {
                    "type": "rendered_content_overflow",
                    "slide": 2,
                    "edges": [{"edge": "right", "mismatch_fraction": 0.12}],
                }
            ]
        },
        "design_quality": {
            "score": 72.5,
            "grade": "acceptable",
            "issues": [{"type": "high_text_density", "slide": 2, "chars": 920, "lines": 18}],
        },
        "visual_quality": {"score": 91.0, "grade": "excellent", "issues": []},
        "fonts": {"issues": [{"type": "expected_font_missing", "font": "Microsoft YaHei"}]},
        "issues": [
            {"type": "shape_out_of_bounds", "slide": 2, "shape": 3},
            {"type": "rendered_content_overflow", "slide": 2, "edges": [{"edge": "right", "mismatch_fraction": 0.12}]},
            {"type": "high_text_density", "slide": 2, "chars": 920, "lines": 18},
            {"type": "expected_font_missing", "font": "Microsoft YaHei"},
        ],
    }

    sections = tool._build_structured_quality_sections(report)

    assert sections["gate"]["status"] == "needs_revision"
    assert sections["gate"]["blocking_issue_count"] == 3
    assert sections["gate"]["affected_slides"] == [2]
    assert sections["metrics"]["design_score"] == 72.5
    assert sections["metrics"]["visual_score"] == 91.0
    assert sections["issue_summary"] == {
        "expected_font_missing": 1,
        "high_text_density": 1,
        "rendered_content_overflow": 1,
        "shape_out_of_bounds": 1,
    }

    issue = sections["structured_issues"][0]
    assert {"id", "type", "category", "severity", "message", "slide", "location", "evidence", "artifacts"}.issubset(issue)
    assert issue["severity"] == "high"
    assert issue["message"]
    assert issue["location"]["shape_index"] == 3
    assert issue["artifacts"]["page_png"] == "/tmp/page-002.png"
    assert "action" not in issue
    assert "recommendation" not in issue


def test_rendered_visual_quality_rejects_sparse_top_left_content(tmp_path: Path):
    image_path = tmp_path / "broken-slide.png"
    image = Image.new("RGB", (1440, 810), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((8, 8, 260, 70), fill="black")
    image.save(image_path)

    report = ValidatePptxTool()._inspect_rendered_visual_quality([image_path])

    issue_types = {issue["type"] for issue in report["issues"]}
    assert "rendered_sparse_or_blank" in issue_types
    assert "rendered_corner_cluster" in issue_types
    assert report["score"] <= 40
    sections = ValidatePptxTool()._build_structured_quality_sections({
        "pages": [{"slide": 1, "png_path": str(image_path)}],
        "issues": report["issues"],
        "visual_quality": report,
    })
    assert sections["gate"]["blocking"] is True


def test_rendered_visual_quality_accepts_deliberately_sparse_cover(tmp_path: Path):
    image_path = tmp_path / "minimal-cover.png"
    image = Image.new("RGB", (1440, 810), "#0b2940")
    draw = ImageDraw.Draw(image)
    # A restrained title cover can legitimately occupy about 2% of the canvas.
    # It must not be treated as a failed/blank render merely for using whitespace.
    draw.rectangle((220, 300, 720, 345), fill="white")
    draw.rectangle((220, 390, 560, 410), fill="#00b5d8")
    image.save(image_path)

    report = ValidatePptxTool()._inspect_rendered_visual_quality([image_path])

    issue_types = {issue["type"] for issue in report["issues"]}
    assert report["slides"][0]["ink_ratio"] > 0.018
    assert "rendered_sparse_or_blank" not in issue_types
    assert not any("接近空白" in item for item in report["recommendations"])


@pytest.mark.asyncio
async def test_ppt_master_shape_plan_ids_are_reported_by_geometry_qa(tmp_path: Path):
    result = await CreatePptxWithPptMasterTool().execute(
        title="Shape ID 定位测试",
        slide_plan=[
            {
                "title": "越界页",
                "shapes": [
                    {
                        "id": "s2_title",
                        "type": "text",
                        "text": "越界标题",
                        "x": 12.8,
                        "y": 0.2,
                        "w": 1.2,
                        "h": 0.4,
                        "font_size": 18,
                    }
                ],
            }
        ],
        output_file=str(tmp_path / "shape_id.pptx"),
        project_dir=str(tmp_path / "shape_id_project"),
        enable_preview=False,
        run_validation=False,
    )

    assert result["success"] is True
    validation = await ValidatePptxTool().execute(
        result["data"]["file_path"],
        render_png=False,
    )

    assert validation["success"] is True
    issues = validation["data"]["structured_issues"]
    out_of_bounds = next(issue for issue in issues if issue["type"] == "shape_out_of_bounds")
    assert out_of_bounds["location"]["shape_id"] == "s2_title"
    assert out_of_bounds["location"]["shape_name"] == "pptm:s2_title"


def test_validate_pptx_toc_like_slide_is_not_flagged_text_only(tmp_path: Path):
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(914400, 914400, 7315200, 914400)
    title.text = "汇报大纲"
    item1 = slide.shapes.add_textbox(1828800, 2286000, 9144000, 365760)
    item1.text = "1. 我们做了什么"
    item2 = slide.shapes.add_textbox(1828800, 2743200, 9144000, 365760)
    item2.text = "2. 怎么做的"
    pptx_path = tmp_path / "toc_like.pptx"
    presentation.save(pptx_path)

    report = ValidatePptxTool()._inspect_design_quality(pptx_path)

    assert not any(issue["type"] == "text_only_slide" for issue in report["issues"])


def test_validate_pptx_rendered_visual_quality_flags_overcrowding(tmp_path: Path):
    image_module = pytest.importorskip("PIL.Image")
    draw_module = pytest.importorskip("PIL.ImageDraw")
    image = image_module.new("RGB", (640, 360), "white")
    draw = draw_module.Draw(image)
    for x in range(8, 620, 20):
        for y in range(8, 340, 20):
            draw.rectangle((x, y, x + 17, y + 17), fill=(20, 20, 20))
    png_path = tmp_path / "dense.png"
    image.save(png_path)

    report = ValidatePptxTool()._inspect_rendered_visual_quality([png_path])

    assert report["enabled"] is True
    assert report["score"] < 100
    assert any(issue["type"] == "rendered_visual_overcrowding" for issue in report["issues"])


def test_validate_pptx_visual_quality_ignores_single_pixel_renderer_seam(tmp_path: Path):
    image_module = pytest.importorskip("PIL.Image")
    draw_module = pytest.importorskip("PIL.ImageDraw")
    image = image_module.new("RGB", (1921, 1080), (10, 37, 63))
    draw = draw_module.Draw(image)
    draw.rectangle((360, 240, 1500, 840), fill=(0, 180, 216))
    draw.line((1920, 0, 1920, 1079), fill="white", width=1)
    png_path = tmp_path / "renderer-seam.png"
    image.save(png_path)

    report = ValidatePptxTool()._inspect_rendered_visual_quality([png_path])

    assert not any(issue["type"] == "rendered_low_margin" for issue in report["issues"])


def test_validate_pptx_visual_quality_still_flags_substantive_edge_content(tmp_path: Path):
    image_module = pytest.importorskip("PIL.Image")
    draw_module = pytest.importorskip("PIL.ImageDraw")
    image = image_module.new("RGB", (640, 360), "white")
    draw = draw_module.Draw(image)
    draw.rectangle((0, 40, 500, 320), fill=(20, 20, 20))
    png_path = tmp_path / "real-edge-content.png"
    image.save(png_path)

    report = ValidatePptxTool()._inspect_rendered_visual_quality([png_path])

    assert any(issue["type"] == "rendered_low_margin" for issue in report["issues"])


def test_ppt_master_quality_gate_returns_rewrite_when_validation_fails():
    gate = CreatePptxWithPptMasterTool()._workflow_quality_gate(
        [
            {"slide": 1, "layout": "cover_statement", "role": "cover"},
            {"slide": 2, "layout": "card_grid", "role": "content"},
        ],
        {
            "success": False,
            "design_quality": {"issues": [{"type": "text_only_slide", "slide": 2}]},
            "overflow_issues": [{"type": "rendered_low_margin", "slide": 3}],
        },
    )

    assert gate["status"] == "rewrite_required"
    assert gate["rewrite_required"] is True
    assert gate["qa_status"] == "needs_revision"
    assert gate["affected_slides"] == [2, 3]
    assert gate["issue_summary"]["validation_failed"] == 1
    assert any(issue["type"] == "validation_failed" for issue in gate["issues"])
    assert any(issue["type"] == "text_only_slide" for issue in gate["issues"])
    assert any(task["slide"] == 2 and task["message"] for task in gate["revision_tasks"])
    assert all("action" not in task for task in gate["revision_tasks"])
    assert any(task["slide"] == 3 and task["priority"] == "high" for task in gate["revision_tasks"])


def test_ppt_master_quality_gate_preserves_structured_issue_context():
    gate = CreatePptxWithPptMasterTool()._workflow_quality_gate(
        [
            {"slide": 1, "layout": "cover_statement", "role": "cover"},
            {"slide": 2, "layout": "agent_shape_plan", "role": "content"},
        ],
        {
            "success": False,
            "structured_issues": [
                {
                    "id": "pptqa-001",
                    "type": "shape_out_of_bounds",
                    "category": "geometry",
                    "severity": "high",
                    "message": "形状边界超出幻灯片画布。",
                    "slide": 2,
                    "location": {"shape_id": "s2_title", "shape_index": 3},
                    "evidence": {"bounds": {"left": 100}},
                    "artifacts": {"page_png": "/tmp/page-002.png"},
                }
            ],
        },
    )

    task = gate["revision_tasks"][0]
    assert task["type"] == "shape_out_of_bounds"
    assert task["priority"] == "high"
    assert task["category"] == "geometry"
    assert task["message"] == "形状边界超出幻灯片画布。"
    assert task["location"]["shape_id"] == "s2_title"
    assert task["evidence"]["bounds"]["left"] == 100
    assert task["artifacts"]["page_png"] == "/tmp/page-002.png"
    assert "action" not in task


def test_ppt_master_quality_gate_distinguishes_qa_failed_from_revision_needed():
    gate = CreatePptxWithPptMasterTool()._workflow_quality_gate(
        [{"slide": 1, "layout": "cover_statement", "role": "cover"}],
        {
            "success": False,
            "issues": [{"type": "validation_error", "message": "render failed"}],
        },
    )

    assert gate["status"] == "qa_failed"
    assert gate["qa_status"] == "qa_failed"
    assert gate["rewrite_required"] is False
    assert gate["revision_tasks"] == []


def test_ppt_master_quality_summary_promotes_revision_loop():
    tool = CreatePptxWithPptMasterTool()
    summary = tool._build_summary(
        output_name="demo.pptx",
        slide_count=5,
        quality_gate={
            "qa_status": "needs_revision",
            "revision_tasks": [
                {"slide": 2, "type": "text_only_slide", "message": "页面主要由文本框组成。"},
                {"slide": 4, "type": "rendered_low_margin", "message": "渲染内容距离页面边缘过近。"},
            ],
        },
    )

    assert "已生成PPT初稿" in summary
    assert "QA发现 2 项需优化任务" in summary
    assert "继续迭代" in summary
