"""Basic PPTX geometry and rendered page checks."""
from __future__ import annotations

import tempfile
from pathlib import Path
from typing import Dict, List

from PIL import Image, ImageStat

EMU_PER_INCH = 914400
PAD_RGB = (200, 200, 200)


def inspect_pptx_geometry(pptx_path: Path) -> Dict[str, object]:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    issues: List[Dict[str, object]] = []
    empty_slides: List[int] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        visible_count = 0
        for shape_index, shape in enumerate(slide.shapes, start=1):
            left = int(getattr(shape, "left", 0))
            top = int(getattr(shape, "top", 0))
            width = int(getattr(shape, "width", 0))
            height = int(getattr(shape, "height", 0))
            if width > 0 and height > 0:
                visible_count += 1
            if left < 0 or top < 0 or left + width > slide_w or top + height > slide_h:
                issues.append(
                    {
                        "type": "shape_out_of_bounds",
                        "slide": slide_index,
                        "shape": shape_index,
                        "bounds": {
                            "left": left,
                            "top": top,
                            "width": width,
                            "height": height,
                        },
                    }
                )

        if visible_count == 0:
            empty_slides.append(slide_index)
            issues.append({"type": "empty_slide", "slide": slide_index})

    return {
        "slide_width": slide_w,
        "slide_height": slide_h,
        "empty_slides": empty_slides,
        "issues": issues,
    }


def inspect_rendered_pages(page_pngs: List[Path], blank_threshold: float = 1.2) -> Dict[str, object]:
    issues: List[Dict[str, object]] = []
    blank_pages: List[int] = []

    for index, path in enumerate(page_pngs, start=1):
        image = Image.open(path).convert("L")
        stat = ImageStat.Stat(image)
        # A nearly all-white rendered page has very low standard deviation.
        stddev = float(stat.stddev[0])
        mean = float(stat.mean[0])
        if mean > 248 and stddev < blank_threshold:
            blank_pages.append(index)
            issues.append(
                {
                    "type": "rendered_blank_page",
                    "slide": index,
                    "mean": round(mean, 2),
                    "stddev": round(stddev, 2),
                }
            )

    return {
        "blank_pages": blank_pages,
        "issues": issues,
    }


def _px_to_emu(px: int, dpi: int) -> int:
    return int(px * EMU_PER_INCH / max(1, dpi))


def _calc_tolerance(dpi: int) -> int:
    if dpi >= 300:
        return 0
    return min(max(round((300 - dpi) / 25), 1), 10)


def _enlarge_deck_with_padding(src: Path, dst: Path, pad_emu: int) -> Dict[str, float]:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Emu

    prs = Presentation(str(src))
    original_w = int(prs.slide_width)
    original_h = int(prs.slide_height)
    new_w = original_w + 2 * pad_emu
    new_h = original_h + 2 * pad_emu
    prs.slide_width = Emu(new_w)
    prs.slide_height = Emu(new_h)

    for slide in prs.slides:
        for shape in list(slide.shapes):
            shape.left = Emu(int(shape.left) + pad_emu)
            shape.top = Emu(int(shape.top) + pad_emu)

        pads = (
            (0, 0, pad_emu, new_h),
            (new_w - pad_emu, 0, pad_emu, new_h),
            (0, 0, new_w, pad_emu),
            (0, new_h - pad_emu, new_w, pad_emu),
        )
        sp_tree = slide.shapes._spTree
        for left, top, width, height in pads:
            pad_shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Emu(left),
                Emu(top),
                Emu(width),
                Emu(height),
            )
            pad_shape.fill.solid()
            pad_shape.fill.fore_color.rgb = RGBColor(*PAD_RGB)
            pad_shape.line.fill.background()
            sp_tree.remove(pad_shape._element)
            sp_tree.insert(2, pad_shape._element)

    prs.save(str(dst))
    return {
        "pad_ratio_w": pad_emu / new_w,
        "pad_ratio_h": pad_emu / new_h,
    }


def _inspect_padding_margins(page_pngs: List[Path], pad_ratio_w: float, pad_ratio_h: float, dpi: int) -> List[Dict[str, object]]:
    try:
        import numpy as np
    except ImportError:
        return [{"type": "rendered_overflow_check_unavailable", "message": "numpy is not installed"}]

    tolerance = _calc_tolerance(dpi)
    pad_colour = np.array(PAD_RGB, dtype=np.uint8)
    issues: List[Dict[str, object]] = []

    for index, path in enumerate(page_pngs, start=1):
        with Image.open(path) as image:
            arr = np.asarray(image.convert("RGB"))

        height, width, _ = arr.shape
        pad_x = max(1, int(width * pad_ratio_w) - 1)
        pad_y = max(1, int(height * pad_ratio_h) - 1)
        margins = {
            "left": arr[:, :pad_x, :],
            "right": arr[:, width - pad_x:, :],
            "top": arr[:pad_y, :, :],
            "bottom": arr[height - pad_y:, :, :],
        }

        dirty_edges = []
        for edge, margin in margins.items():
            diff = np.abs(margin.astype("int16") - pad_colour.astype("int16"))
            matches = np.all(diff <= tolerance, axis=-1)
            mismatch_fraction = 1.0 - (np.count_nonzero(matches) / matches.size)
            max_mismatch = 0.03 if dpi < 200 else 0.02
            if mismatch_fraction > max_mismatch:
                dirty_edges.append({"edge": edge, "mismatch_fraction": round(float(mismatch_fraction), 4)})

        if dirty_edges:
            issues.append(
                {
                    "type": "rendered_content_overflow",
                    "slide": index,
                    "edges": dirty_edges,
                }
            )

    return issues


def inspect_rendered_overflow(pptx_path: Path, qa_dir: Path, dpi: int = 144, pad_px: int = 100) -> Dict[str, object]:
    from app.tools.office.slides_qa.render_pptx import render_deck

    overflow_dir = qa_dir / "overflow"
    overflow_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="pptx_overflow_") as tmp_name:
        enlarged_path = Path(tmp_name) / "enlarged.pptx"
        pad_info = _enlarge_deck_with_padding(pptx_path, enlarged_path, _px_to_emu(pad_px, dpi))
        render_result = render_deck(enlarged_path, overflow_dir, dpi=dpi)
        page_pngs = [Path(path) for path in render_result.get("page_pngs", [])]
        issues = _inspect_padding_margins(
            page_pngs,
            float(pad_info["pad_ratio_w"]),
            float(pad_info["pad_ratio_h"]),
            dpi,
        )

    return {
        "enabled": True,
        "pad_px": pad_px,
        "render": render_result,
        "issues": issues,
    }
