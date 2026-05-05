"""
Create a PPTX from an existing template by replacing analyzed slots.
"""
from __future__ import annotations

import shutil
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional
from urllib.parse import quote

import structlog

from app.tools.base.tool_interface import LLMTool, ToolCategory
from app.tools.office.analyze_pptx_template_tool import AnalyzePptxTemplateTool

logger = structlog.get_logger()


class CreatePptxFromTemplateTool(LLMTool):
    def __init__(self):
        super().__init__(
            name="create_pptx_from_template",
            description="基于PPTX模板和slot映射替换文本/表格/图片，生成新的可编辑PPTX。",
            category=ToolCategory.QUERY,
            version="1.0.0",
            requires_context=False,
        )
        self.working_dir = Path.cwd().parent
        self.default_output_dir = self.working_dir / "backend" / "backend_data_registry" / "presentations"

    async def execute(
        self,
        template_path: str,
        replacements: Any,
        output_file: Optional[str] = None,
        run_validation: bool = False,
        quality: str = "draft",
        **kwargs,
    ) -> Dict[str, Any]:
        try:
            from pptx import Presentation
        except ImportError:
            return {
                "success": False,
                "data": {"error": "python-pptx 未安装，请在 suyuan 环境安装 python-pptx"},
                "summary": "模板生成PPT失败：缺少 python-pptx",
            }

        try:
            template = self._resolve_path(template_path)
            if not template.exists():
                return {
                    "success": False,
                    "data": {"error": f"模板文件不存在: {template}"},
                    "summary": "模板生成PPT失败：模板文件不存在",
                }
            if template.suffix.lower() != ".pptx":
                return {
                    "success": False,
                    "data": {"error": f"只支持 .pptx 模板，当前格式: {template.suffix}"},
                    "summary": "模板生成PPT失败：格式不支持",
                }

            output_path = self._resolve_output_file(output_file, template)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(template, output_path)

            analysis_result = await AnalyzePptxTemplateTool().execute(
                str(output_path),
                include_layouts=False,
                write_report=False,
            )
            if not analysis_result.get("success"):
                return {
                    "success": False,
                    "data": {"error": analysis_result.get("data", {}).get("error", "模板分析失败")},
                    "summary": "模板生成PPT失败：模板分析失败",
                }

            slot_map = self._build_slot_map(analysis_result["data"])
            normalized_replacements = self._normalize_replacements(replacements)
            prs = Presentation(str(output_path))
            operations: List[Dict[str, Any]] = []
            warnings: List[Dict[str, Any]] = []

            for slot_id, value in normalized_replacements.items():
                slot = slot_map.get(slot_id)
                if not slot:
                    warnings.append({"slot_id": slot_id, "warning": "slot_id 不存在，已跳过"})
                    continue
                try:
                    operation = self._apply_replacement(prs, slot, value)
                    operations.append(operation)
                except Exception as exc:
                    logger.warning("template_slot_replace_failed", slot_id=slot_id, error=str(exc))
                    warnings.append({"slot_id": slot_id, "warning": str(exc)})

            prs.save(str(output_path))

            from config.settings import settings

            api_base = settings.backend_host.rstrip("/")
            encoded_path = quote(str(output_path))
            result_data: Dict[str, Any] = {
                "file_path": str(output_path),
                "output_file": str(output_path),
                "file_name": output_path.name,
                "template_path": str(template),
                "operation_count": len(operations),
                "warning_count": len(warnings),
                "operations": operations,
                "warnings": warnings,
                "doc_url": f"{api_base}/api/utility/file/{encoded_path}",
                "doc_download_filename": output_path.name,
            }

            quality_mode = str(quality or "draft").lower()
            if quality_mode not in {"draft", "standard", "strict"}:
                quality_mode = "draft"
            if run_validation or quality_mode in {"standard", "strict"}:
                try:
                    from app.tools.office.validate_pptx_tool import ValidatePptxTool

                    validation = await ValidatePptxTool().execute(
                        str(output_path),
                        render_overflow_check=quality_mode == "strict" or run_validation,
                    )
                    result_data["validation"] = validation.get("data")
                except Exception as validation_error:
                    logger.warning("template_pptx_validation_failed", error=str(validation_error))
                    result_data["validation_error"] = str(validation_error)

            return {
                "success": True,
                "data": result_data,
                "summary": f"已基于模板生成PPT：{output_path.name}，替换 {len(operations)} 个槽位",
            }
        except Exception as e:
            logger.error("create_pptx_from_template_failed", template_path=template_path, error=str(e), exc_info=True)
            return {
                "success": False,
                "data": {"error": str(e)},
                "summary": f"模板生成PPT失败：{str(e)[:80]}",
            }

    def _apply_replacement(self, prs, slot: Dict[str, Any], value: Any) -> Dict[str, Any]:
        slide = prs.slides[int(slot["slide"]) - 1]
        shape = slide.shapes[int(slot["shape_index"]) - 1]
        replacement_type = self._replacement_type(slot, value)

        if replacement_type == "image":
            image_path = self._image_path(value)
            if not image_path or not image_path.exists():
                raise ValueError(f"图片不存在: {image_path}")
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            self._remove_shape(shape)
            slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        elif replacement_type == "table":
            rows = value.get("rows") if isinstance(value, dict) else value
            if not isinstance(rows, list):
                raise ValueError("table 替换值必须是二维数组或 {rows: ...}")
            if not getattr(shape, "has_table", False):
                raise ValueError("目标槽位不是表格")
            self._replace_table(shape, rows)
        else:
            text = self._text_value(value)
            if not getattr(shape, "has_text_frame", False):
                raise ValueError("目标槽位不是文本框")
            shape.text = text

        return {
            "slot_id": slot["slot_id"],
            "slide": slot["slide"],
            "shape_index": slot["shape_index"],
            "type": replacement_type,
        }

    def _replace_table(self, shape, rows: List[Any]) -> None:
        table = shape.table
        row_count = len(table.rows)
        col_count = len(table.columns)
        for r_idx in range(row_count):
            for c_idx in range(col_count):
                value = ""
                if r_idx < len(rows) and isinstance(rows[r_idx], list) and c_idx < len(rows[r_idx]):
                    value = rows[r_idx][c_idx]
                table.cell(r_idx, c_idx).text = "" if value is None else str(value)

    def _remove_shape(self, shape) -> None:
        element = shape._element
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)

    def _build_slot_map(self, analysis: Dict[str, Any]) -> Dict[str, Dict[str, Any]]:
        slot_map: Dict[str, Dict[str, Any]] = {}
        for slide in analysis.get("slides", []):
            for slot in slide.get("replaceable_slots", []):
                slot_map[slot["slot_id"]] = slot
        return slot_map

    def _normalize_replacements(self, replacements: Any) -> Dict[str, Any]:
        if isinstance(replacements, dict):
            if "slots" in replacements and isinstance(replacements["slots"], list):
                return self._normalize_replacements(replacements["slots"])
            return replacements
        if isinstance(replacements, list):
            normalized: Dict[str, Any] = {}
            for item in replacements:
                if not isinstance(item, dict):
                    continue
                slot_id = item.get("slot_id")
                if not slot_id:
                    continue
                normalized[str(slot_id)] = item.get("value", item.get("text", item.get("rows", item.get("path"))))
            return normalized
        return {}

    def _replacement_type(self, slot: Dict[str, Any], value: Any) -> str:
        if isinstance(value, dict):
            explicit_type = value.get("type")
            if explicit_type in {"text", "table", "image"}:
                return explicit_type
            if value.get("path"):
                return "image"
            if value.get("rows"):
                return "table"
        if isinstance(value, list):
            return "table"
        slot_type = slot.get("type")
        return "table" if slot_type == "table" else "text"

    def _image_path(self, value: Any) -> Optional[Path]:
        raw = value.get("path") if isinstance(value, dict) else value
        if not raw:
            return None
        return self._resolve_path(str(raw))

    def _text_value(self, value: Any) -> str:
        if isinstance(value, dict):
            value = value.get("text", value.get("value", ""))
        if isinstance(value, list):
            return "\n".join("" if item is None else str(item) for item in value)
        return "" if value is None else str(value)

    def _resolve_path(self, path: str) -> Path:
        file_path = Path(path)
        if not file_path.is_absolute():
            file_path = self.working_dir / file_path
        return file_path.resolve()

    def _resolve_output_file(self, output_file: Optional[str], template: Path) -> Path:
        if output_file:
            path = Path(output_file)
            if not path.is_absolute():
                path = self.working_dir / path
        else:
            path = self.default_output_dir / f"{template.stem}_filled_{uuid.uuid4().hex[:8]}.pptx"
        if path.suffix.lower() != ".pptx":
            path = path.with_suffix(".pptx")
        return path.resolve()

    def get_function_schema(self) -> Dict[str, Any]:
        return {
            "name": "create_pptx_from_template",
            "description": "复制PPTX模板并按 analyze_pptx_template 产出的 slot_id 替换文本、表格或图片，生成新的可编辑PPTX。",
            "parameters": {
                "type": "object",
                "properties": {
                    "template_path": {"type": "string", "description": "PPTX模板路径"},
                    "replacements": {
                        "type": "object",
                        "description": "槽位替换映射。推荐格式: {\"s001_slot001\": \"新文本\", \"s002_slot006\": [[\"A\",\"B\"]], \"s003_slot009\": {\"type\":\"image\",\"path\":\"...\"}}",
                    },
                    "output_file": {"type": "string", "description": "输出PPTX路径，可选"},
                    "run_validation": {"type": "boolean", "description": "是否生成后调用validate_pptx", "default": False},
                    "quality": {
                        "type": "string",
                        "description": "质量模式：draft/standard/strict",
                        "default": "draft",
                        "enum": ["draft", "standard", "strict"],
                    },
                },
                "required": ["template_path", "replacements"],
            },
        }

    def is_available(self) -> bool:
        return True


tool = CreatePptxFromTemplateTool()
