"""
Read PPTX files into structured slide content for Agent use.
"""
import hashlib
from pathlib import Path
from typing import Any, Dict, List, Optional

import structlog

from app.tools.base.tool_interface import LLMTool, ToolCategory

logger = structlog.get_logger()


def get_pdf_converter():
    try:
        from app.services.pdf_converter import pdf_converter
        return pdf_converter
    except ImportError:
        logger.warning("pdf_converter_not_available")
        return None


class ReadPptxTool(LLMTool):
    def __init__(self):
        super().__init__(
            name="read_pptx",
            description="读取PPTX演示文稿，提取每页文本、表格、图片信息、备注和基础元数据。",
            category=ToolCategory.QUERY,
            version="1.0.0",
            requires_context=False,
        )
        self.working_dir = Path.cwd().parent

    async def execute(
        self,
        path: str,
        include_notes: bool = True,
        include_images: bool = True,
        export_images: bool = False,
        image_output_dir: Optional[str] = None,
        max_slides: Optional[int] = None,
        enable_preview: bool = True,
        **kwargs,
    ) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return {
                "success": False,
                "data": {"error": "python-pptx 未安装，请在 suyuan 环境安装 python-pptx"},
                "summary": "读取PPT失败：缺少 python-pptx",
            }

        try:
            file_path = self._resolve_path(path)
            if not file_path.exists():
                return {
                    "success": False,
                    "data": {"error": f"文件不存在: {file_path}"},
                    "summary": "读取PPT失败：文件不存在",
                }
            if file_path.suffix.lower() != ".pptx":
                return {
                    "success": False,
                    "data": {"error": f"只支持 .pptx 文件，当前格式: {file_path.suffix}"},
                    "summary": "读取PPT失败：格式不支持",
                }

            prs = Presentation(str(file_path))
            slide_total = len(prs.slides)
            limit = slide_total if max_slides is None else max(0, min(max_slides, slide_total))
            slides: List[Dict[str, Any]] = []
            export_dir = self._resolve_image_output_dir(file_path, image_output_dir) if export_images else None

            for idx, slide in enumerate(list(prs.slides)[:limit], start=1):
                slide_data: Dict[str, Any] = {
                    "index": idx,
                    "title": "",
                    "texts": [],
                    "tables": [],
                    "images": [],
                    "notes": "",
                }

                if slide.shapes.title and getattr(slide.shapes.title, "text", None):
                    slide_data["title"] = slide.shapes.title.text.strip()

                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        text = shape.text.strip()
                        if text:
                            slide_data["texts"].append(text)
                            if not slide_data["title"]:
                                slide_data["title"] = text.splitlines()[0].strip()

                    if getattr(shape, "has_table", False):
                        rows = []
                        for row in shape.table.rows:
                            rows.append([cell.text.strip() for cell in row.cells])
                        slide_data["tables"].append(rows)

                    if include_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        image_data = {
                            "filename": image.filename,
                            "content_type": image.content_type,
                            "width": int(shape.width),
                            "height": int(shape.height),
                        }
                        if export_dir:
                            image_data["path"] = str(self._export_image(export_dir, idx, image))
                        slide_data["images"].append(image_data)

                if include_notes:
                    try:
                        notes_frame = slide.notes_slide.notes_text_frame
                        if notes_frame:
                            slide_data["notes"] = notes_frame.text.strip()
                    except Exception:
                        slide_data["notes"] = ""

                slides.append(slide_data)

            result_data: Dict[str, Any] = {
                "file_path": str(file_path),
                "file_name": file_path.name,
                "slide_count": slide_total,
                "returned_slide_count": len(slides),
                "slide_width": int(prs.slide_width),
                "slide_height": int(prs.slide_height),
                "slides": slides,
            }
            if export_dir:
                result_data["image_output_dir"] = str(export_dir)

            if enable_preview:
                try:
                    converter = get_pdf_converter()
                    if converter:
                        result_data["pdf_preview"] = await converter.convert_to_pdf(str(file_path))
                except Exception as preview_error:
                    logger.warning("read_pptx_preview_failed", error=str(preview_error))

            return {
                "success": True,
                "data": result_data,
                "summary": f"已读取 {file_path.name}，共 {slide_total} 页，返回 {len(slides)} 页内容",
            }
        except Exception as e:
            logger.error("read_pptx_failed", path=path, error=str(e), exc_info=True)
            return {
                "success": False,
                "data": {"error": str(e)},
                "summary": f"读取PPT失败：{str(e)[:80]}",
            }

    def _resolve_path(self, path: str) -> Path:
        file_path = Path(path)
        if not file_path.is_absolute():
            file_path = self.working_dir / file_path
        return file_path.resolve()

    def _resolve_image_output_dir(self, file_path: Path, image_output_dir: Optional[str]) -> Path:
        if image_output_dir:
            output_dir = Path(image_output_dir)
            if not output_dir.is_absolute():
                output_dir = self.working_dir / output_dir
        else:
            output_dir = (
                self.working_dir
                / "backend"
                / "backend_data_registry"
                / "presentations"
                / "assets"
                / file_path.stem
            )
        output_dir.mkdir(parents=True, exist_ok=True)
        return output_dir.resolve()

    def _export_image(self, output_dir: Path, slide_index: int, image) -> Path:
        ext = Path(image.filename or "").suffix
        if not ext:
            content_type = getattr(image, "content_type", "")
            ext = {
                "image/png": ".png",
                "image/jpeg": ".jpg",
                "image/gif": ".gif",
                "image/bmp": ".bmp",
                "image/svg+xml": ".svg",
            }.get(content_type, ".bin")
        digest = hashlib.sha1(image.blob).hexdigest()[:10]
        output_path = output_dir / f"slide_{slide_index:03d}_{digest}{ext.lower()}"
        if not output_path.exists():
            output_path.write_bytes(image.blob)
        return output_path

    def get_function_schema(self) -> Dict[str, Any]:
        return {
            "name": "read_pptx",
            "description": "读取PPTX，提取幻灯片文本、表格、图片元数据、备注，并可生成PDF预览。",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "PPTX文件路径"},
                    "include_notes": {"type": "boolean", "description": "是否读取备注", "default": True},
                    "include_images": {"type": "boolean", "description": "是否返回图片元数据", "default": True},
                    "export_images": {"type": "boolean", "description": "是否导出PPT内嵌图片到文件", "default": False},
                    "image_output_dir": {"type": "string", "description": "图片导出目录，可选"},
                    "max_slides": {"type": "integer", "description": "最多读取多少页，可选"},
                    "enable_preview": {"type": "boolean", "description": "是否生成PDF预览", "default": True},
                },
                "required": ["path"],
            },
        }

    def is_available(self) -> bool:
        return True


tool = ReadPptxTool()
