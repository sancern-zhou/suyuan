"""
Analyze PPTX templates into a structured slot map for later template-based generation/editing.
"""
from __future__ import annotations

import hashlib
import json
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional

import structlog

from app.tools.base.tool_interface import LLMTool, ToolCategory

logger = structlog.get_logger()


EMU_PER_INCH = 914400


class AnalyzePptxTemplateTool(LLMTool):
    def __init__(self):
        super().__init__(
            name="analyze_pptx_template",
            description="分析PPTX模板结构，输出页面、版式、占位符、文本/图片/表格槽位和模板地图。",
            category=ToolCategory.QUERY,
            version="1.0.0",
            requires_context=False,
        )
        self.working_dir = Path.cwd().parent
        self.default_output_dir = self.working_dir / "backend" / "backend_data_registry" / "presentations" / "template_analysis"

    async def execute(
        self,
        path: str,
        max_slides: Optional[int] = None,
        include_layouts: bool = True,
        include_hidden_text: bool = False,
        write_report: bool = True,
        output_dir: Optional[str] = None,
        **kwargs,
    ) -> Dict[str, Any]:
        try:
            from pptx import Presentation
        except ImportError:
            return {
                "success": False,
                "data": {"error": "python-pptx 未安装，请在 suyuan 环境安装 python-pptx"},
                "summary": "PPT模板分析失败：缺少 python-pptx",
            }

        try:
            pptx_path = self._resolve_path(path)
            if not pptx_path.exists():
                return {
                    "success": False,
                    "data": {"error": f"文件不存在: {pptx_path}"},
                    "summary": "PPT模板分析失败：文件不存在",
                }
            if pptx_path.suffix.lower() != ".pptx":
                return {
                    "success": False,
                    "data": {"error": f"只支持 .pptx 文件，当前格式: {pptx_path.suffix}"},
                    "summary": "PPT模板分析失败：格式不支持",
                }

            prs = Presentation(str(pptx_path))
            slide_total = len(prs.slides)
            limit = slide_total if max_slides is None else max(0, min(max_slides, slide_total))
            slide_w = int(prs.slide_width)
            slide_h = int(prs.slide_height)

            slides: List[Dict[str, Any]] = []
            slot_index = 1
            for slide_index, slide in enumerate(list(prs.slides)[:limit], start=1):
                slide_data, next_slot_index = self._analyze_slide(
                    slide=slide,
                    slide_index=slide_index,
                    slot_start=slot_index,
                    slide_w=slide_w,
                    slide_h=slide_h,
                    include_hidden_text=include_hidden_text,
                )
                slot_index = next_slot_index
                slides.append(slide_data)

            layouts = self._analyze_layouts(prs) if include_layouts else []
            report = {
                "file_path": str(pptx_path),
                "file_name": pptx_path.name,
                "slide_count": slide_total,
                "returned_slide_count": len(slides),
                "slide_width": slide_w,
                "slide_height": slide_h,
                "slide_width_in": round(slide_w / EMU_PER_INCH, 3),
                "slide_height_in": round(slide_h / EMU_PER_INCH, 3),
                "layouts": layouts,
                "slides": slides,
                "summary": self._build_summary(slides),
            }

            if write_report:
                report_path = self._write_report(report, output_dir, pptx_path)
                report["report_path"] = str(report_path)

            return {
                "success": True,
                "data": report,
                "summary": (
                    f"已分析PPT模板：{pptx_path.name}，共 {slide_total} 页，"
                    f"返回 {len(slides)} 页，识别 {sum(len(s['replaceable_slots']) for s in slides)} 个可替换槽位"
                ),
            }
        except Exception as e:
            logger.error("analyze_pptx_template_failed", path=path, error=str(e), exc_info=True)
            return {
                "success": False,
                "data": {"error": str(e)},
                "summary": f"PPT模板分析失败：{str(e)[:80]}",
            }

    def _analyze_slide(
        self,
        slide,
        slide_index: int,
        slot_start: int,
        slide_w: int,
        slide_h: int,
        include_hidden_text: bool,
    ) -> tuple[Dict[str, Any], int]:
        shape_items: List[Dict[str, Any]] = []
        placeholders: List[Dict[str, Any]] = []
        text_boxes: List[Dict[str, Any]] = []
        pictures: List[Dict[str, Any]] = []
        tables: List[Dict[str, Any]] = []
        charts: List[Dict[str, Any]] = []
        replaceable_slots: List[Dict[str, Any]] = []
        slot_index = slot_start

        for shape_index, shape in enumerate(slide.shapes, start=1):
            shape_data = self._shape_data(shape, shape_index, slide_w, slide_h)
            text = self._shape_text(shape)
            if text and (include_hidden_text or not self._looks_hidden(shape_data)):
                shape_data["text"] = text
            shape_items.append(shape_data)

            if shape_data["is_placeholder"]:
                placeholder = self._placeholder_data(shape, shape_data, text)
                placeholders.append(placeholder)
                replaceable_slots.append(self._slot_data(slot_index, slide_index, shape_index, shape_data, text, "placeholder"))
                slot_index += 1
            elif text:
                text_box = {
                    "shape_index": shape_index,
                    "text": text,
                    "bounds": shape_data["bounds"],
                    "role": self._infer_text_role(text, shape_data),
                }
                text_boxes.append(text_box)
                if not self._is_non_content_text(text, shape_data):
                    replaceable_slots.append(self._slot_data(slot_index, slide_index, shape_index, shape_data, text, "text"))
                    slot_index += 1

            if shape_data["kind"] == "picture":
                picture = self._picture_data(shape, shape_index, shape_data)
                pictures.append(picture)
                replaceable_slots.append(self._slot_data(slot_index, slide_index, shape_index, shape_data, "", "image"))
                slot_index += 1

            if getattr(shape, "has_table", False):
                table = self._table_data(shape, shape_index, shape_data)
                tables.append(table)
                replaceable_slots.append(self._slot_data(slot_index, slide_index, shape_index, shape_data, "", "table"))
                slot_index += 1

            if getattr(shape, "has_chart", False):
                chart = self._chart_data(shape, shape_index, shape_data)
                charts.append(chart)
                replaceable_slots.append(self._slot_data(slot_index, slide_index, shape_index, shape_data, "", "chart"))
                slot_index += 1

        title = self._slide_title(slide, text_boxes, placeholders)
        for slot in replaceable_slots:
            if slot.get("current_text") and slot["current_text"].splitlines()[0].strip() == title:
                slot["role"] = "title"
        slide_data = {
            "index": slide_index,
            "layout": {
                "name": getattr(slide.slide_layout, "name", ""),
                "index": self._layout_index(slide),
            },
            "title": title,
            "classification": self._classify_slide(title, text_boxes, pictures, tables, charts),
            "shape_count": len(shape_items),
            "placeholders": placeholders,
            "text_boxes": text_boxes,
            "pictures": pictures,
            "tables": tables,
            "charts": charts,
            "replaceable_slots": replaceable_slots,
            "shapes": shape_items,
        }
        return slide_data, slot_index

    def _shape_data(self, shape, shape_index: int, slide_w: int, slide_h: int) -> Dict[str, Any]:
        left = int(getattr(shape, "left", 0) or 0)
        top = int(getattr(shape, "top", 0) or 0)
        width = int(getattr(shape, "width", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)
        return {
            "shape_index": shape_index,
            "name": getattr(shape, "name", ""),
            "kind": self._shape_kind(shape),
            "shape_type": self._enum_name(getattr(shape, "shape_type", "")),
            "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
            "has_text": bool(getattr(shape, "has_text_frame", False)),
            "has_table": bool(getattr(shape, "has_table", False)),
            "has_chart": bool(getattr(shape, "has_chart", False)),
            "bounds": {
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "x_in": round(left / EMU_PER_INCH, 3),
                "y_in": round(top / EMU_PER_INCH, 3),
                "w_in": round(width / EMU_PER_INCH, 3),
                "h_in": round(height / EMU_PER_INCH, 3),
                "relative": {
                    "x": round(left / slide_w, 4) if slide_w else 0,
                    "y": round(top / slide_h, 4) if slide_h else 0,
                    "w": round(width / slide_w, 4) if slide_w else 0,
                    "h": round(height / slide_h, 4) if slide_h else 0,
                },
            },
        }

    def _shape_kind(self, shape) -> str:
        if getattr(shape, "has_table", False):
            return "table"
        if getattr(shape, "has_chart", False):
            return "chart"
        shape_type = str(self._enum_name(getattr(shape, "shape_type", ""))).lower()
        if "picture" in shape_type:
            return "picture"
        if getattr(shape, "has_text_frame", False):
            return "text"
        if "group" in shape_type:
            return "group"
        if "line" in shape_type:
            return "line"
        return "shape"

    def _shape_text(self, shape) -> str:
        if not getattr(shape, "has_text_frame", False):
            return ""
        return (getattr(shape, "text", "") or "").strip()

    def _placeholder_data(self, shape, shape_data: Dict[str, Any], text: str) -> Dict[str, Any]:
        placeholder_format = getattr(shape, "placeholder_format", None)
        return {
            "shape_index": shape_data["shape_index"],
            "name": shape_data["name"],
            "type": self._enum_name(getattr(placeholder_format, "type", "")),
            "idx": getattr(placeholder_format, "idx", None),
            "text": text,
            "bounds": shape_data["bounds"],
            "role": self._infer_placeholder_role(shape, text),
        }

    def _picture_data(self, shape, shape_index: int, shape_data: Dict[str, Any]) -> Dict[str, Any]:
        image = getattr(shape, "image", None)
        blob = getattr(image, "blob", b"") if image else b""
        return {
            "shape_index": shape_index,
            "name": shape_data["name"],
            "filename": getattr(image, "filename", "") if image else "",
            "content_type": getattr(image, "content_type", "") if image else "",
            "sha1": hashlib.sha1(blob).hexdigest()[:12] if blob else "",
            "bounds": shape_data["bounds"],
        }

    def _table_data(self, shape, shape_index: int, shape_data: Dict[str, Any]) -> Dict[str, Any]:
        rows = []
        table = shape.table
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        return {
            "shape_index": shape_index,
            "rows": len(table.rows),
            "columns": len(table.columns),
            "preview": rows[:5],
            "bounds": shape_data["bounds"],
        }

    def _chart_data(self, shape, shape_index: int, shape_data: Dict[str, Any]) -> Dict[str, Any]:
        chart = shape.chart
        return {
            "shape_index": shape_index,
            "chart_type": self._enum_name(getattr(chart, "chart_type", "")),
            "has_title": bool(getattr(chart, "has_title", False)),
            "bounds": shape_data["bounds"],
        }

    def _slot_data(
        self,
        slot_index: int,
        slide_index: int,
        shape_index: int,
        shape_data: Dict[str, Any],
        text: str,
        slot_type: str,
    ) -> Dict[str, Any]:
        slot_id = f"s{slide_index:03d}_slot{slot_index:03d}"
        return {
            "slot_id": slot_id,
            "slide": slide_index,
            "shape_index": shape_index,
            "type": slot_type,
            "role": self._infer_text_role(text, shape_data) if slot_type in {"text", "placeholder"} else slot_type,
            "current_text": text,
            "bounds": shape_data["bounds"],
            "suggested_content": self._suggested_content(slot_type, text, shape_data),
        }

    def _analyze_layouts(self, prs) -> List[Dict[str, Any]]:
        layouts = []
        for index, layout in enumerate(prs.slide_layouts):
            placeholders = []
            for shape in layout.placeholders:
                placeholder_format = getattr(shape, "placeholder_format", None)
                placeholders.append(
                    {
                        "name": getattr(shape, "name", ""),
                        "type": self._enum_name(getattr(placeholder_format, "type", "")),
                        "idx": getattr(placeholder_format, "idx", None),
                        "bounds": {
                            "x_in": round(int(getattr(shape, "left", 0) or 0) / EMU_PER_INCH, 3),
                            "y_in": round(int(getattr(shape, "top", 0) or 0) / EMU_PER_INCH, 3),
                            "w_in": round(int(getattr(shape, "width", 0) or 0) / EMU_PER_INCH, 3),
                            "h_in": round(int(getattr(shape, "height", 0) or 0) / EMU_PER_INCH, 3),
                        },
                    }
                )
            layouts.append(
                {
                    "index": index,
                    "name": getattr(layout, "name", ""),
                    "placeholder_count": len(placeholders),
                    "placeholders": placeholders,
                }
            )
        return layouts

    def _build_summary(self, slides: List[Dict[str, Any]]) -> Dict[str, Any]:
        classifications: Dict[str, int] = {}
        for slide in slides:
            kind = slide["classification"]
            classifications[kind] = classifications.get(kind, 0) + 1
        return {
            "classification_counts": classifications,
            "replaceable_slot_count": sum(len(slide["replaceable_slots"]) for slide in slides),
            "placeholder_count": sum(len(slide["placeholders"]) for slide in slides),
            "picture_count": sum(len(slide["pictures"]) for slide in slides),
            "table_count": sum(len(slide["tables"]) for slide in slides),
            "chart_count": sum(len(slide["charts"]) for slide in slides),
        }

    def _slide_title(self, slide, text_boxes: List[Dict[str, Any]], placeholders: List[Dict[str, Any]]) -> str:
        if slide.shapes.title and getattr(slide.shapes.title, "text", None):
            title = slide.shapes.title.text.strip()
            if title:
                return title
        for placeholder in placeholders:
            if placeholder["role"] == "title" and placeholder["text"]:
                return placeholder["text"].splitlines()[0].strip()
        if text_boxes:
            return text_boxes[0]["text"].splitlines()[0].strip()
        return ""

    def _classify_slide(
        self,
        title: str,
        text_boxes: List[Dict[str, Any]],
        pictures: List[Dict[str, Any]],
        tables: List[Dict[str, Any]],
        charts: List[Dict[str, Any]],
    ) -> str:
        lower_title = title.lower()
        text_count = len(text_boxes)
        if "目录" in title or "contents" in lower_title or "agenda" in lower_title:
            return "toc"
        if "总结" in title or "小结" in title or "summary" in lower_title or "conclusion" in lower_title:
            return "summary"
        if charts:
            return "data_visualization"
        if tables:
            return "table"
        if pictures and text_count >= 1:
            return "image_text"
        if pictures:
            return "image_showcase"
        if text_count <= 1 and title:
            return "cover_or_section"
        if text_count >= 4:
            return "multi_item_content"
        return "content"

    def _infer_placeholder_role(self, shape, text: str) -> str:
        placeholder_format = getattr(shape, "placeholder_format", None)
        placeholder_type = str(self._enum_name(getattr(placeholder_format, "type", ""))).lower()
        if "title" in placeholder_type or "ctr_title" in placeholder_type:
            return "title"
        if "subtitle" in placeholder_type:
            return "subtitle"
        if "picture" in placeholder_type:
            return "image"
        if "chart" in placeholder_type:
            return "chart"
        if "table" in placeholder_type:
            return "table"
        if text:
            return self._infer_text_role(text, {"is_placeholder": True})
        return "content"

    def _infer_text_role(self, text: str, shape_data: Dict[str, Any]) -> str:
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        bounds = shape_data.get("bounds", {})
        h_in = bounds.get("h_in", 0)
        y_in = bounds.get("y_in", 0)
        if y_in < 1.2 and len(lines) <= 2:
            return "title"
        if len(lines) == 1 and h_in <= 0.6:
            return "label"
        if len(lines) >= 3:
            return "body_or_list"
        return "body"

    def _suggested_content(self, slot_type: str, text: str, shape_data: Dict[str, Any]) -> str:
        if slot_type == "image":
            return "图片路径或图片资源"
        if slot_type == "table":
            return "二维表格数据"
        if slot_type == "chart":
            return "图表数据"
        role = self._infer_text_role(text, shape_data)
        if role == "title":
            return "短标题，建议不超过一行"
        if role == "label":
            return "短标签或指标名"
        return "正文、列表项或说明文本"

    def _looks_hidden(self, shape_data: Dict[str, Any]) -> bool:
        bounds = shape_data.get("bounds", {})
        return bounds.get("width", 0) <= 0 or bounds.get("height", 0) <= 0

    def _is_non_content_text(self, text: str, shape_data: Dict[str, Any]) -> bool:
        bounds = shape_data.get("bounds", {})
        relative = bounds.get("relative", {})
        normalized = text.strip()
        if normalized.isdigit() and relative.get("x", 0) > 0.85 and relative.get("y", 0) > 0.9:
            return True
        return False

    def _layout_index(self, slide) -> Optional[int]:
        try:
            layouts = slide.part.package.presentation_part.presentation.slide_layouts
            for index, layout in enumerate(layouts):
                if layout == slide.slide_layout:
                    return index
        except Exception:
            return None
        return None

    def _enum_name(self, value: Any) -> str:
        name = getattr(value, "name", None)
        if name:
            return str(name)
        return str(value)

    def _resolve_path(self, path: str) -> Path:
        file_path = Path(path)
        if not file_path.is_absolute():
            file_path = self.working_dir / file_path
        return file_path.resolve()

    def _resolve_output_dir(self, output_dir: Optional[str], pptx_path: Path) -> Path:
        if output_dir:
            path = Path(output_dir)
            if not path.is_absolute():
                path = self.working_dir / path
            return path.resolve()
        return (self.default_output_dir / f"{pptx_path.stem}_{uuid.uuid4().hex[:8]}").resolve()

    def _write_report(self, report: Dict[str, Any], output_dir: Optional[str], pptx_path: Path) -> Path:
        report_dir = self._resolve_output_dir(output_dir, pptx_path)
        report_dir.mkdir(parents=True, exist_ok=True)
        report_path = report_dir / "template_map.json"
        report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
        return report_path

    def get_function_schema(self) -> Dict[str, Any]:
        return {
            "name": "analyze_pptx_template",
            "description": "分析PPTX模板结构，返回页面分类、版式、占位符、文本/图片/表格/图表槽位和可替换slot地图。",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "PPTX模板文件路径"},
                    "max_slides": {"type": "integer", "description": "最多分析多少页，可选"},
                    "include_layouts": {"type": "boolean", "description": "是否返回母版版式占位符信息", "default": True},
                    "include_hidden_text": {"type": "boolean", "description": "是否包含疑似隐藏文本", "default": False},
                    "write_report": {"type": "boolean", "description": "是否写出template_map.json", "default": True},
                    "output_dir": {"type": "string", "description": "分析报告输出目录，可选"},
                },
                "required": ["path"],
            },
        }

    def is_available(self) -> bool:
        return True


tool = AnalyzePptxTemplateTool()
