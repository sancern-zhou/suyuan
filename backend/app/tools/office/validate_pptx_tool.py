"""
Validate PPTX deliverables by rendering and running lightweight QA checks.
"""
from __future__ import annotations

import json
import hashlib
import re
import uuid
from collections import Counter
from pathlib import Path
from typing import Any, Dict, List, Optional

import structlog

from app.tools.base.tool_interface import LLMTool, ToolCategory
from app.tools.resource_declarations import derivative_file, primary_file
from app.tools.office.slides_qa.create_montage import create_montage
from app.tools.office.slides_qa.detect_fonts import detect_pdf_fonts
from app.tools.office.slides_qa.detect_overflow import (
    inspect_pptx_geometry,
    inspect_rendered_overflow,
    inspect_rendered_pages,
)
from app.tools.office.slides_qa.render_pptx import render_deck
from app.utils.path_config import BACKEND_ROOT, get_data_registry, resolve_agent_path

logger = structlog.get_logger()


def validation_output_resources(pptx_path: Path, outputs: List[Path]) -> List[Dict[str, Any]]:
    """Declare a presentation and all QA derivatives as one resource group."""
    validation_slot = hashlib.sha256(str(pptx_path).encode("utf-8")).hexdigest()[:16]
    group_key = f"presentation:{validation_slot}"
    primary = primary_file(
        pptx_path,
        group_key=group_key,
        tool_name="validate_pptx",
        renderer="presentation",
        capabilities=("preview", "download", "edit"),
    )
    primary["resource_key"] = "pptx"
    members = [primary]
    for output in outputs:
        if not output.is_file():
            continue
        suffix = output.suffix.lower().lstrip(".") or "file"
        is_page = output.name.startswith("page-")
        relation = "attachment" if is_page or suffix == "json" else "preview"
        renderer = (
            "pdf"
            if suffix == "pdf"
            else "image"
            if suffix in {"png", "jpg", "jpeg", "svg"}
            else "file"
        )
        derivative = derivative_file(
            output,
            group_key=group_key,
            parent_key="pptx",
            tool_name="validate_pptx",
            relation=relation,
            renderer=renderer,
        )
        if output.name == "montage.png":
            derivative["resource_key"] = "montage"
        elif suffix == "pdf":
            derivative["resource_key"] = "pdf"
        elif is_page:
            derivative["resource_key"] = output.stem
        else:
            derivative["resource_key"] = output.name
        members.append(derivative)
    return members


class ValidatePptxTool(LLMTool):
    def __init__(self):
        super().__init__(
            name="validate_pptx",
            description=(
                "渲染PPTX并执行基础交付检查：PDF/PNG预览、montage、空页/越界/字体检测。\n\n"
                "⚠️ 使用前请先阅读：app/tools/office/PPT操作指南.md"
            ),
            category=ToolCategory.QUERY,
            version="1.0.0",
            requires_context=False,
        )
        self.working_dir = BACKEND_ROOT
        self.default_qa_root = get_data_registry() / "presentations" / "qa"

    async def execute(
        self,
        path: str,
        output_dir: Optional[str] = None,
        expected_fonts: Optional[List[str]] = None,
        render_png: bool = True,
        create_overview: bool = True,
        render_overflow_check: bool = True,
        dpi: int = 144,
        **kwargs,
    ) -> Dict[str, Any]:
        try:
            pptx_path = self._resolve_path(path)
            if not pptx_path.exists():
                return {
                    "success": False,
                    "data": {"error": f"文件不存在: {pptx_path}"},
                    "summary": "PPT验证失败：文件不存在",
                }
            if pptx_path.suffix.lower() != ".pptx":
                return {
                    "success": False,
                    "data": {"error": f"只支持 .pptx 文件，当前格式: {pptx_path.suffix}"},
                    "summary": "PPT验证失败：格式不支持",
                }

            qa_dir = self._resolve_output_dir(output_dir, pptx_path)
            qa_dir.mkdir(parents=True, exist_ok=True)

            geometry = inspect_pptx_geometry(pptx_path)
            render_result: Dict[str, Any] = {}
            rendered_checks: Dict[str, Any] = {"issues": [], "blank_pages": []}
            overflow_checks: Dict[str, Any] = {"enabled": False, "issues": []}
            font_checks: Dict[str, Any] = {"fonts": [], "issues": [], "missing_expected_fonts": []}
            design_quality = self._inspect_design_quality(pptx_path)
            visual_quality: Dict[str, Any] = {"enabled": False, "issues": []}
            montage_path = None

            if render_png:
                try:
                    render_result = render_deck(pptx_path, qa_dir, dpi=dpi)
                    page_pngs = [Path(path) for path in render_result.get("page_pngs", [])]
                    rendered_checks = inspect_rendered_pages(page_pngs)
                    visual_quality = self._inspect_rendered_visual_quality(page_pngs)
                    if render_overflow_check:
                        overflow_checks = inspect_rendered_overflow(pptx_path, qa_dir, dpi=dpi)
                    if create_overview:
                        montage_path = create_montage(page_pngs, qa_dir / "montage.png")
                    pdf_path = render_result.get("pdf_path")
                    if pdf_path:
                        font_checks = detect_pdf_fonts(Path(str(pdf_path)), expected_fonts=expected_fonts)
                except Exception as render_error:
                    rendered_checks = {
                        "issues": [
                            {
                                "type": "rendered_visual_qa_unavailable",
                                "message": str(render_error),
                            }
                        ],
                        "blank_pages": [],
                    }
                    visual_quality = {
                        "enabled": False,
                        "score": None,
                        "issues": rendered_checks["issues"],
                        "recommendations": ["安装 PyMuPDF/fitz 后可启用渲染级视觉质量检查。"],
                    }
            page_pngs = [Path(path) for path in render_result.get("page_pngs", [])]
            pages = [
                {
                    "slide": index,
                    "page_number": index,
                    "file_name": page_path.name,
                    "png_path": str(page_path),
                }
                for index, page_path in enumerate(page_pngs, start=1)
            ]

            issues = []
            issues.extend(geometry.get("issues", []))
            issues.extend(rendered_checks.get("issues", []))
            issues.extend(overflow_checks.get("issues", []))
            issues.extend(font_checks.get("issues", []))
            issues.extend(design_quality.get("issues", []))
            issues.extend(visual_quality.get("issues", []))

            report = {
                "success": len(issues) == 0,
                "pptx_path": str(pptx_path),
                "qa_dir": str(qa_dir),
                "render": render_result,
                "montage_path": str(montage_path) if montage_path else None,
                "page_index_base": 1,
                "pages": pages,
                "usage_notes": [
                    "pages 中的 slide/page_number 与 page-XXX.png 文件名均为 1-based。",
                    "montage_path 仅用于全局总览；修复或核查指定页时应读取 pages 中对应 slide 的 png_path。",
                ],
                "geometry": geometry,
                "rendered_pages": rendered_checks,
                "rendered_overflow": overflow_checks,
                "fonts": font_checks,
                "design_quality": design_quality,
                "visual_quality": visual_quality,
                "issues": issues,
                "issue_count": len(issues),
            }
            report.update(self._build_structured_quality_sections(report))

            report_path = qa_dir / "report.json"
            report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
            report["report_path"] = str(report_path)

            compact_outputs = [report_path]
            if montage_path:
                compact_outputs.append(montage_path)
            pdf_path = render_result.get("pdf_path") if isinstance(render_result, dict) else None
            if pdf_path:
                compact_outputs.append(Path(str(pdf_path)))
            compact_outputs.extend(page_pngs)

            summary = (
                f"PPT验证完成：{pptx_path.name}，发现 {len(issues)} 个问题"
                if issues
                else f"PPT验证通过：{pptx_path.name}"
            )
            return {
                "success": True,
                "data": report,
                "resources": validation_output_resources(pptx_path, compact_outputs),
                "summary": summary,
            }
        except Exception as e:
            logger.error("validate_pptx_failed", path=path, error=str(e), exc_info=True)
            return {
                "success": False,
                "data": {"error": str(e)},
                "summary": f"PPT验证失败：{str(e)[:80]}",
            }

    def _resolve_path(self, path: str) -> Path:
        return resolve_agent_path(path)

    def _resolve_output_dir(self, output_dir: Optional[str], pptx_path: Path) -> Path:
        if output_dir:
            return resolve_agent_path(output_dir)
        return (self.default_qa_root / f"{pptx_path.stem}_{uuid.uuid4().hex[:8]}").resolve()

    def _build_structured_quality_sections(self, report: Dict[str, Any]) -> Dict[str, Any]:
        raw_issues = [issue for issue in report.get("issues", []) if isinstance(issue, dict)]
        structured_issues = [
            self._structure_issue(issue, index, report)
            for index, issue in enumerate(raw_issues, start=1)
        ]
        issue_summary = dict(sorted(Counter(issue["type"] for issue in structured_issues).items()))
        affected_slides = sorted(
            {
                issue["slide"]
                for issue in structured_issues
                if isinstance(issue.get("slide"), int)
            }
        )
        blocking_issue_count = sum(1 for issue in structured_issues if issue["severity"] == "high")
        qa_failed = bool(structured_issues) and all(
            issue["type"] in {"validation_error", "rendered_visual_qa_unavailable"}
            for issue in structured_issues
        )
        status = "qa_failed" if qa_failed else ("needs_revision" if structured_issues else "passed")
        return {
            "metrics": self._build_quality_metrics(report),
            "structured_issues": structured_issues,
            "issue_summary": issue_summary,
            "gate": {
                "status": status,
                "passed": status == "passed",
                "blocking": status == "qa_failed" or blocking_issue_count > 0,
                "blocking_issue_count": blocking_issue_count,
                "issue_count": len(structured_issues),
                "affected_slides": affected_slides,
            },
        }

    def _build_quality_metrics(self, report: Dict[str, Any]) -> Dict[str, Any]:
        design_quality = report.get("design_quality") if isinstance(report.get("design_quality"), dict) else {}
        visual_quality = report.get("visual_quality") if isinstance(report.get("visual_quality"), dict) else {}
        rendered_pages = report.get("rendered_pages") if isinstance(report.get("rendered_pages"), dict) else {}
        geometry = report.get("geometry") if isinstance(report.get("geometry"), dict) else {}
        fonts = report.get("fonts") if isinstance(report.get("fonts"), dict) else {}
        pages = report.get("pages") if isinstance(report.get("pages"), list) else []
        return {
            "slide_count": len(pages),
            "design_score": design_quality.get("score"),
            "design_grade": design_quality.get("grade"),
            "visual_score": visual_quality.get("score"),
            "visual_grade": visual_quality.get("grade"),
            "blank_pages": rendered_pages.get("blank_pages", []),
            "empty_slides": geometry.get("empty_slides", []),
            "missing_expected_fonts": fonts.get("missing_expected_fonts", []),
        }

    def _structure_issue(self, issue: Dict[str, Any], index: int, report: Dict[str, Any]) -> Dict[str, Any]:
        issue_type = str(issue.get("type") or "unknown")
        slide = issue.get("slide")
        if not isinstance(slide, int):
            slide = None
        structured = {
            "id": f"pptqa-{index:03d}",
            "type": issue_type,
            "category": self._issue_category(issue_type),
            "severity": self._issue_severity(issue_type),
            "message": self._issue_message(issue_type, issue),
            "slide": slide,
            "location": self._issue_location(issue),
            "evidence": self._issue_evidence(issue),
            "artifacts": self._issue_artifacts(slide, report),
            "raw_issue": issue,
        }
        return structured

    def _issue_category(self, issue_type: str) -> str:
        if issue_type in {"shape_out_of_bounds", "empty_slide"}:
            return "geometry"
        if issue_type.startswith("rendered_"):
            return "rendered_visual"
        if issue_type in {"high_text_density", "moderate_text_density", "text_only_slide", "tiny_text", "weak_title_hierarchy", "low_typographic_contrast"}:
            return "design_structure"
        if issue_type in {"expected_font_missing", "font_detection_unavailable"}:
            return "font"
        if issue_type in {"validation_error", "validation_failed"}:
            return "qa_runtime"
        return "general"

    def _issue_severity(self, issue_type: str) -> str:
        high = {
            "shape_out_of_bounds",
            "empty_slide",
            "rendered_blank_page",
            "rendered_content_overflow",
            "rendered_low_margin",
            "rendered_visual_overcrowding",
            "rendered_sparse_or_blank",
            "rendered_nearly_blank",
            "rendered_corner_cluster",
            "high_text_density",
            "validation_error",
            "rendered_visual_qa_unavailable",
        }
        medium = {
            "text_only_slide",
            "tiny_text",
            "expected_font_missing",
            "repeated_layout_pattern",
            "rendered_dense_composition",
        }
        if issue_type in high:
            return "high"
        if issue_type in medium:
            return "medium"
        return "low"

    def _issue_message(self, issue_type: str, issue: Dict[str, Any]) -> str:
        messages = {
            "shape_out_of_bounds": "形状边界超出幻灯片画布。",
            "empty_slide": "页面没有可见形状。",
            "rendered_blank_page": "渲染结果接近空白。",
            "rendered_content_overflow": "渲染后页面边缘检测到非背景内容。",
            "rendered_low_margin": "渲染内容距离页面边缘过近。",
            "rendered_visual_overcrowding": "渲染视觉占用过高。",
            "rendered_sparse_or_blank": "渲染结果内容过少或接近空白。",
            "rendered_nearly_blank": "渲染结果接近完全空白。",
            "rendered_corner_cluster": "渲染内容异常集中在页面左上角。",
            "high_text_density": "页面文字行数或字符数超过高密度阈值。",
            "moderate_text_density": "页面文字行数或字符数接近密集阈值。",
            "text_only_slide": "页面主要由文本框组成，缺少图表、图片、表格或形状等视觉元素。",
            "tiny_text": "页面存在小于可读阈值的字号。",
            "expected_font_missing": "PDF 字体检测未发现期望字体。",
            "validation_error": "PPT 验证流程执行异常。",
            "rendered_visual_qa_unavailable": "渲染级视觉 QA 不可用。",
        }
        message = messages.get(issue_type, "检测到 PPT 质量问题。")
        detail = issue.get("message")
        if detail:
            return f"{message} {detail}"
        return message

    def _issue_location(self, issue: Dict[str, Any]) -> Dict[str, Any]:
        location: Dict[str, Any] = {}
        if isinstance(issue.get("shape"), int):
            location["shape_index"] = issue["shape"]
        if issue.get("shape_id"):
            location["shape_id"] = issue["shape_id"]
        if issue.get("shape_name"):
            location["shape_name"] = issue["shape_name"]
        if isinstance(issue.get("bounds"), dict):
            location["bounds"] = issue["bounds"]
        if isinstance(issue.get("edges"), list):
            location["edges"] = issue["edges"]
        return location

    def _issue_evidence(self, issue: Dict[str, Any]) -> Dict[str, Any]:
        excluded = {"type", "slide", "shape", "bounds", "message"}
        return {
            key: value
            for key, value in issue.items()
            if key not in excluded
        }

    def _issue_artifacts(self, slide: Optional[int], report: Dict[str, Any]) -> Dict[str, Any]:
        artifacts: Dict[str, Any] = {}
        if report.get("montage_path"):
            artifacts["montage"] = report["montage_path"]
        if not isinstance(slide, int):
            return artifacts
        for page in report.get("pages", []):
            if isinstance(page, dict) and page.get("slide") == slide and page.get("png_path"):
                artifacts["page_png"] = page["png_path"]
                break
        return artifacts

    def _inspect_rendered_visual_quality(self, page_pngs: List[Path]) -> Dict[str, Any]:
        try:
            from PIL import Image, ImageChops
        except ImportError:
            return {
                "enabled": False,
                "score": None,
                "issues": [{"type": "visual_quality_unavailable", "message": "Pillow 未安装"}],
            }

        slides: List[Dict[str, Any]] = []
        issues: List[Dict[str, Any]] = []
        for slide_index, path in enumerate(page_pngs, start=1):
            with Image.open(path) as image:
                rgb = image.convert("RGB")
                sample = rgb.resize((320, 180))

            bg = self._estimate_background_color(sample)
            bg_image = Image.new("RGB", sample.size, bg)
            diff = ImageChops.difference(sample, bg_image).convert("L")
            threshold = 18
            changed = diff.point(lambda value: 255 if value > threshold else 0)
            width, height = sample.size
            # LibreOffice may emit a one-pixel white seam on the right edge
            # (for example a 1921px raster for a 16:9 slide). Ignore only the
            # outer sample pixel; substantive edge content remains visible in
            # the adjacent pixels and is still flagged.
            changed.paste(0, (0, 0, width, 1))
            changed.paste(0, (0, height - 1, width, height))
            changed.paste(0, (0, 0, 1, height))
            changed.paste(0, (width - 1, 0, width, height))
            bbox = changed.getbbox()
            changed_pixels = sum(1 for value in changed.getdata() if value)
            ink_ratio = changed_pixels / float(width * height)

            score = 100
            slide_issues: List[Dict[str, Any]] = []
            margin_ratio = None
            if bbox:
                left, top, right, bottom = bbox
                margin_ratio = min(left / width, top / height, (width - right) / width, (height - bottom) / height)
                if margin_ratio < 0.015 and ink_ratio > 0.08:
                    score -= 18
                    slide_issues.append(
                        {
                            "type": "rendered_low_margin",
                            "slide": slide_index,
                            "margin_ratio": round(margin_ratio, 4),
                        }
                    )
                if left / width < 0.05 and top / height < 0.08 and right / width < 0.55 and bottom / height < 0.5:
                    score -= 25
                    slide_issues.append(
                        {
                            "type": "rendered_corner_cluster",
                            "slide": slide_index,
                            "content_bbox_ratio": [
                                round(left / width, 4), round(top / height, 4),
                                round(right / width, 4), round(bottom / height, 4),
                            ],
                        }
                    )
            else:
                score -= 35
                slide_issues.append({"type": "rendered_nearly_blank", "slide": slide_index, "ink_ratio": 0})

            if ink_ratio > 0.58:
                score -= 30
                slide_issues.append(
                    {"type": "rendered_visual_overcrowding", "slide": slide_index, "ink_ratio": round(ink_ratio, 4)}
                )
            elif ink_ratio > 0.48:
                score -= 16
                slide_issues.append(
                    {"type": "rendered_dense_composition", "slide": slide_index, "ink_ratio": round(ink_ratio, 4)}
                )
            # Keep the threshold low enough to catch failed renders while
            # allowing intentional title covers with generous negative space.
            # Real blank/near-blank pages and tiny corner clusters remain
            # blocking via this threshold and the bbox checks above.
            elif ink_ratio < 0.018:
                score -= 45
                slide_issues.append(
                    {"type": "rendered_sparse_or_blank", "slide": slide_index, "ink_ratio": round(ink_ratio, 4)}
                )

            score = max(0, score)
            issues.extend(
                issue
                for issue in slide_issues
                if issue["type"] in {
                    "rendered_visual_overcrowding",
                    "rendered_low_margin",
                    "rendered_nearly_blank",
                    "rendered_sparse_or_blank",
                    "rendered_corner_cluster",
                }
            )
            slides.append(
                {
                    "slide": slide_index,
                    "score": score,
                    "ink_ratio": round(ink_ratio, 4),
                    "margin_ratio": round(margin_ratio, 4) if margin_ratio is not None else None,
                    "background_rgb": bg,
                    "issues": slide_issues,
                }
            )

        overall = round(sum(item["score"] for item in slides) / len(slides), 1) if slides else 0
        return {
            "enabled": True,
            "score": overall,
            "grade": self._design_grade(overall),
            "slides": slides,
            "issues": issues,
            "recommendations": self._visual_recommendations(slides),
        }

    def _estimate_background_color(self, image: Any) -> tuple[int, int, int]:
        width, height = image.size
        points = [
            image.getpixel((0, 0)),
            image.getpixel((width - 1, 0)),
            image.getpixel((0, height - 1)),
            image.getpixel((width - 1, height - 1)),
            image.getpixel((width // 2, 0)),
            image.getpixel((width // 2, height - 1)),
        ]
        return tuple(sorted(channel)[len(channel) // 2] for channel in zip(*points))

    def _visual_recommendations(self, slides: List[Dict[str, Any]]) -> List[str]:
        recommendations: List[str] = []
        if any(item["ink_ratio"] > 0.48 for item in slides):
            recommendations.append("渲染结果视觉占用过高，建议减少同页元素、拆页或把正文改成图表/流程。")
        if any(
            item["margin_ratio"] is not None
            and item["margin_ratio"] < 0.015
            and item["ink_ratio"] > 0.08
            for item in slides
        ):
            recommendations.append("渲染结果贴边明显，建议增加页边距并缩小主内容区域。")
        if any(item["ink_ratio"] < 0.018 for item in slides):
            recommendations.append("渲染结果接近空白，建议检查图片、图表或文字是否成功输出。")
        if any(any(issue["type"] == "rendered_corner_cluster" for issue in item["issues"]) for item in slides):
            recommendations.append("内容集中在左上角，建议检查 CSS/主题是否加载，并重新建立页面视觉层级。")
        return recommendations

    def _inspect_design_quality(self, pptx_path: Path) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return {
                "enabled": False,
                "score": None,
                "issues": [{"type": "design_quality_unavailable", "message": "python-pptx 未安装"}],
            }

        prs = Presentation(str(pptx_path))
        slide_scores: List[Dict[str, Any]] = []
        issues: List[Dict[str, Any]] = []
        previous_signature = None
        repeated_layouts = 0

        for slide_index, slide in enumerate(prs.slides, start=1):
            text_chars = 0
            text_lines = 0
            text_boxes = 0
            visual_shapes = 0
            pictures = 0
            tables = 0
            charts = 0
            font_sizes: List[float] = []
            shape_signature = []
            slide_text_parts: List[str] = []

            for shape in slide.shapes:
                left = round(int(getattr(shape, "left", 0) or 0) / 914400, 1)
                top = round(int(getattr(shape, "top", 0) or 0) / 914400, 1)
                width = round(int(getattr(shape, "width", 0) or 0) / 914400, 1)
                height = round(int(getattr(shape, "height", 0) or 0) / 914400, 1)
                shape_signature.append((left, top, width, height))

                if getattr(shape, "has_text_frame", False):
                    text = (getattr(shape, "text", "") or "").strip()
                    if text:
                        slide_text_parts.append(text)
                        text_boxes += 1
                        text_chars += len(text)
                        text_lines += max(1, len([line for line in text.splitlines() if line.strip()]))
                    if len(text) > 2:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                size = getattr(run.font, "size", None)
                                if size is not None:
                                    font_sizes.append(round(size.pt, 1))
                if getattr(shape, "has_table", False):
                    tables += 1
                    visual_shapes += 1
                if getattr(shape, "has_chart", False):
                    charts += 1
                    visual_shapes += 1
                shape_type = getattr(shape, "shape_type", None)
                if shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pictures += 1
                    visual_shapes += 1
                elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # Native rectangles, circles, cards, and other drawn shapes
                    # expose an (often empty) text frame in python-pptx. They are
                    # still visual structure and must not be mistaken for text boxes.
                    visual_shapes += 1
                elif not getattr(shape, "has_text_frame", False):
                    visual_shapes += 1

            signature = tuple(shape_signature[:8])
            if previous_signature == signature and slide_index > 1:
                repeated_layouts += 1
            previous_signature = signature

            score = 100
            slide_issues: List[Dict[str, Any]] = []
            # A card/dashboard layout commonly uses many short text boxes. Count
            # it as high density only when line count and actual copy volume agree.
            if text_chars > 900 or (text_lines > 30 and text_chars > 600):
                score -= 25
                slide_issues.append(
                    {"type": "high_text_density", "slide": slide_index, "chars": text_chars, "lines": text_lines}
                )
            elif text_chars > 620 or (text_lines > 20 and text_chars > 320):
                score -= 15
                slide_issues.append(
                    {"type": "moderate_text_density", "slide": slide_index, "chars": text_chars, "lines": text_lines}
                )
            slide_text = "\n".join(slide_text_parts)
            is_toc_like = self._is_toc_like_slide(slide_text)
            if visual_shapes == 0 and text_boxes >= 2 and slide_index > 1 and not is_toc_like:
                score -= 20
                slide_issues.append({"type": "text_only_slide", "slide": slide_index})
            if font_sizes:
                max_font = max(font_sizes)
                min_font = min(font_sizes)
                if max_font < 24 and slide_index > 1:
                    score -= 10
                    slide_issues.append({"type": "weak_title_hierarchy", "slide": slide_index, "max_font": max_font})
                if min_font < 8:
                    score -= 10
                    slide_issues.append({"type": "tiny_text", "slide": slide_index, "min_font": min_font})
                if max_font - min_font < 8 and text_boxes >= 2:
                    score -= 8
                    slide_issues.append({"type": "low_typographic_contrast", "slide": slide_index})

            score = max(0, score)
            issues.extend(issue for issue in slide_issues if issue["type"] in {"high_text_density", "text_only_slide", "tiny_text"})
            slide_scores.append(
                {
                    "slide": slide_index,
                    "score": score,
                    "text_chars": text_chars,
                    "text_lines": text_lines,
                    "text_boxes": text_boxes,
                    "visual_shapes": visual_shapes,
                    "pictures": pictures,
                    "tables": tables,
                    "charts": charts,
                    "issues": slide_issues,
                }
            )

        if repeated_layouts >= max(2, len(prs.slides) // 3):
            issues.append({"type": "repeated_layout_pattern", "count": repeated_layouts})

        overall = round(sum(item["score"] for item in slide_scores) / len(slide_scores), 1) if slide_scores else 0
        if repeated_layouts:
            overall = max(0, overall - min(12, repeated_layouts * 3))
        return {
            "enabled": True,
            "score": overall,
            "grade": self._design_grade(overall),
            "repeated_layouts": repeated_layouts,
            "slides": slide_scores,
            "issues": issues,
            "recommendations": self._design_recommendations(slide_scores, repeated_layouts),
        }

    def _is_toc_like_slide(self, text: str) -> bool:
        if not text:
            return False
        patterns = [
            r"目录",
            r"大纲",
            r"汇报大纲",
            r"agenda",
            r"toc",
        ]
        return any(re.search(pattern, text, flags=re.IGNORECASE) for pattern in patterns)

    def _design_grade(self, score: float) -> str:
        if score >= 90:
            return "excellent"
        if score >= 78:
            return "good"
        if score >= 65:
            return "acceptable"
        return "needs_improvement"

    def _design_recommendations(self, slide_scores: List[Dict[str, Any]], repeated_layouts: int) -> List[str]:
        recommendations: List[str] = []
        if any(item["text_chars"] > 620 or item["text_lines"] > 10 for item in slide_scores):
            recommendations.append("拆分高文字密度页面，或改为主结论+卡片/图表结构。")
        if any(item["visual_shapes"] == 0 and item["text_boxes"] >= 2 and item["slide"] > 1 for item in slide_scores):
            recommendations.append("为纯文字页面补充图表、流程、指标卡或领域示意组件。")
        if repeated_layouts:
            recommendations.append("连续页面版式过于接近，建议交替使用双栏、卡片、图表和时间线布局。")
        if any(any(issue["type"] == "weak_title_hierarchy" for issue in item["issues"]) for item in slide_scores):
            recommendations.append("提高标题字号或降低正文强调，增强标题/正文层级。")
        return recommendations

    def get_function_schema(self) -> Dict[str, Any]:
        return {
            "name": "validate_pptx",
            "description": "渲染PPTX为PDF/PNG并执行基础QA检查，返回montage、report.json、字体信息和设计质量评分。",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "PPTX文件路径"},
                    "output_dir": {"type": "string", "description": "QA输出目录，可选"},
                    "expected_fonts": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "期望字体名称列表，可选",
                    },
                    "render_png": {"type": "boolean", "description": "是否渲染PNG页面", "default": True},
                    "create_overview": {"type": "boolean", "description": "是否生成montage总览图", "default": True},
                    "render_overflow_check": {
                        "type": "boolean",
                        "description": "是否执行渲染级溢出检测，默认开启",
                        "default": True,
                    },
                    "dpi": {"type": "integer", "description": "PNG渲染DPI，默认144"},
                },
                "required": ["path"],
            },
        }

    def is_available(self) -> bool:
        return True


tool = ValidatePptxTool()
