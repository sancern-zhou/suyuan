"""
Edit existing PPTX files with common safe operations.
"""
from __future__ import annotations

import shutil
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional
from urllib.parse import quote

import structlog

from app.tools.base.tool_interface import LLMTool, ToolCategory
from app.tools.office.analyze_pptx_template_tool import AnalyzePptxTemplateTool
from app.tools.office.create_pptx_from_template_tool import CreatePptxFromTemplateTool

logger = structlog.get_logger()


class EditPptxTool(LLMTool):
    def __init__(self):
        super().__init__(
            name="edit_pptx",
            description="编辑现有PPTX：全局替换文本、按slot替换内容、删除页、重排页，并可执行PPT质量验证。",
            category=ToolCategory.QUERY,
            version="1.0.0",
            requires_context=False,
        )
        self.working_dir = Path.cwd().parent
        self.default_output_dir = self.working_dir / "backend" / "backend_data_registry" / "presentations"

    async def execute(
        self,
        path: str,
        operations: Optional[List[Dict[str, Any]]] = None,
        output_file: Optional[str] = None,
        replacements: Optional[Any] = None,
        run_validation: bool = False,
        quality: str = "draft",
        **kwargs,
    ) -> Dict[str, Any]:
        try:
            from pptx import Presentation
        except ImportError:
            return {
                "success": False,
                "data": {"error": "python-pptx 未安装，请在 suyuan 环境安装 python-pptx"},
                "summary": "编辑PPT失败：缺少 python-pptx",
            }

        try:
            source = self._resolve_path(path)
            if not source.exists():
                return {
                    "success": False,
                    "data": {"error": f"文件不存在: {source}"},
                    "summary": "编辑PPT失败：文件不存在",
                }
            if source.suffix.lower() != ".pptx":
                return {
                    "success": False,
                    "data": {"error": f"只支持 .pptx 文件，当前格式: {source.suffix}"},
                    "summary": "编辑PPT失败：格式不支持",
                }

            output_path = self._resolve_output_file(output_file, source)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            if source.resolve() != output_path.resolve():
                shutil.copy2(source, output_path)

            prs = Presentation(str(output_path))
            normalized_ops = self._normalize_operations(operations, replacements, kwargs)
            applied: List[Dict[str, Any]] = []
            warnings: List[Dict[str, Any]] = []

            for index, operation in enumerate(normalized_ops, start=1):
                op_type = str(operation.get("type") or operation.get("operation") or "").strip()
                try:
                    if op_type == "replace_text":
                        applied.append(self._replace_text(prs, operation))
                    elif op_type == "replace_slot":
                        slot_operations, slot_warnings = await self._replace_slots(prs, str(output_path), operation)
                        applied.extend(slot_operations)
                        warnings.extend({"operation_index": index, **warning} for warning in slot_warnings)
                    elif op_type == "delete_slides":
                        applied.append(self._delete_slides(prs, operation))
                    elif op_type == "reorder_slides":
                        applied.append(self._reorder_slides(prs, operation))
                    else:
                        warnings.append({"operation_index": index, "warning": f"未知操作类型: {op_type or '<empty>'}"})
                except Exception as exc:
                    logger.warning("edit_pptx_operation_failed", operation=operation, error=str(exc))
                    warnings.append({"operation_index": index, "operation": operation, "warning": str(exc)})

            prs.save(str(output_path))

            from config.settings import settings

            api_base = settings.backend_host.rstrip("/")
            encoded_path = quote(str(output_path))
            result_data: Dict[str, Any] = {
                "file_path": str(output_path),
                "output_file": str(output_path),
                "file_name": output_path.name,
                "source_path": str(source),
                "slide_count": len(prs.slides),
                "operation_count": len(applied),
                "warning_count": len(warnings),
                "operations": applied,
                "warnings": warnings,
                "doc_url": f"{api_base}/api/utility/file/{encoded_path}",
                "doc_download_filename": output_path.name,
            }

            quality_mode = str(quality or "draft").lower()
            if quality_mode not in {"draft", "standard", "strict"}:
                quality_mode = "draft"
            if run_validation or quality_mode in {"standard", "strict"}:
                try:
                    from app.tools.office.validate_pptx_tool import ValidatePptxTool

                    validation = await ValidatePptxTool().execute(
                        str(output_path),
                        render_overflow_check=quality_mode == "strict" or run_validation,
                    )
                    result_data["validation"] = validation.get("data")
                except Exception as validation_error:
                    logger.warning("edit_pptx_validation_failed", error=str(validation_error))
                    result_data["validation_error"] = str(validation_error)

            return {
                "success": True,
                "data": result_data,
                "summary": f"已编辑PPT：{output_path.name}，完成 {len(applied)} 个操作",
            }
        except Exception as e:
            logger.error("edit_pptx_failed", path=path, error=str(e), exc_info=True)
            return {
                "success": False,
                "data": {"error": str(e)},
                "summary": f"编辑PPT失败：{str(e)[:80]}",
            }

    def _replace_text(self, prs, operation: Dict[str, Any]) -> Dict[str, Any]:
        old = operation.get("old_text", operation.get("search", operation.get("find")))
        new = operation.get("new_text", operation.get("replace", operation.get("value", "")))
        if old is None:
            raise ValueError("replace_text 需要 old_text/search")
        old_text = str(old)
        new_text = "" if new is None else str(new)
        slide_filter = self._slide_filter(operation.get("slides"))
        replacements = 0

        for slide_index, slide in enumerate(prs.slides, start=1):
            if slide_filter and slide_index not in slide_filter:
                continue
            for shape in slide.shapes:
                replacements += self._replace_in_shape(shape, old_text, new_text)

        return {
            "type": "replace_text",
            "old_text": old_text,
            "new_text": new_text,
            "replacement_count": replacements,
        }

    async def _replace_slots(self, prs, output_path: str, operation: Dict[str, Any]) -> tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
        replacements = operation.get("replacements")
        if replacements is None:
            slot_id = operation.get("slot_id")
            if not slot_id:
                raise ValueError("replace_slot 需要 slot_id 或 replacements")
            replacements = {str(slot_id): operation.get("value", operation.get("text", operation.get("rows", operation.get("path"))))}

        analysis_result = await AnalyzePptxTemplateTool().execute(
            output_path,
            include_layouts=False,
            write_report=False,
        )
        if not analysis_result.get("success"):
            raise ValueError(analysis_result.get("data", {}).get("error", "PPT slot 分析失败"))

        helper = CreatePptxFromTemplateTool()
        slot_map = helper._build_slot_map(analysis_result["data"])
        normalized = helper._normalize_replacements(replacements)
        applied: List[Dict[str, Any]] = []
        warnings: List[Dict[str, Any]] = []

        for slot_id, value in normalized.items():
            slot = slot_map.get(slot_id)
            if not slot:
                warnings.append({"slot_id": slot_id, "warning": "slot_id 不存在，已跳过"})
                continue
            result = helper._apply_replacement(prs, slot, value)
            result["type"] = f"replace_slot:{result['type']}"
            applied.append(result)

        return applied, warnings

    def _delete_slides(self, prs, operation: Dict[str, Any]) -> Dict[str, Any]:
        indexes = operation.get("slides", operation.get("indexes", operation.get("slide_indexes")))
        if not isinstance(indexes, list) or not indexes:
            raise ValueError("delete_slides 需要 slides 数组，页码从 1 开始")
        slide_count = len(prs.slides)
        unique_indexes = sorted({int(item) for item in indexes}, reverse=True)
        for slide_index in unique_indexes:
            if slide_index < 1 or slide_index > slide_count:
                raise ValueError(f"删除页码越界: {slide_index}，当前共 {slide_count} 页")

        sld_id_lst = prs.slides._sldIdLst
        slide_ids = list(sld_id_lst)
        for slide_index in unique_indexes:
            sld_id_lst.remove(slide_ids[slide_index - 1])

        return {"type": "delete_slides", "deleted_slides": sorted(unique_indexes), "slide_count": len(prs.slides)}

    def _reorder_slides(self, prs, operation: Dict[str, Any]) -> Dict[str, Any]:
        order = operation.get("order", operation.get("slides"))
        if not isinstance(order, list) or not order:
            raise ValueError("reorder_slides 需要 order 数组，页码从 1 开始")
        slide_count = len(prs.slides)
        normalized_order = [int(item) for item in order]
        if sorted(normalized_order) != list(range(1, slide_count + 1)):
            raise ValueError(f"reorder_slides 必须包含 1..{slide_count} 的完整页码且不能重复")

        sld_id_lst = prs.slides._sldIdLst
        slide_ids = list(sld_id_lst)
        for slide_id in slide_ids:
            sld_id_lst.remove(slide_id)
        for slide_index in normalized_order:
            sld_id_lst.append(slide_ids[slide_index - 1])

        return {"type": "reorder_slides", "order": normalized_order, "slide_count": len(prs.slides)}

    def _replace_in_shape(self, shape, old_text: str, new_text: str) -> int:
        replacements = 0
        if getattr(shape, "has_text_frame", False):
            replacements += self._replace_in_text_frame(shape.text_frame, old_text, new_text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    replacements += self._replace_in_text_frame(cell.text_frame, old_text, new_text)
        return replacements

    def _replace_in_text_frame(self, text_frame, old_text: str, new_text: str) -> int:
        replacements = 0
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    count = run.text.count(old_text)
                    run.text = run.text.replace(old_text, new_text)
                    replacements += count
        if replacements == 0 and old_text in text_frame.text:
            count = text_frame.text.count(old_text)
            text_frame.text = text_frame.text.replace(old_text, new_text)
            replacements += count
        return replacements

    def _slide_filter(self, slides: Any) -> set:
        if not slides:
            return set()
        if not isinstance(slides, list):
            slides = [slides]
        return {int(slide) for slide in slides}

    def _normalize_operations(
        self,
        operations: Optional[List[Dict[str, Any]]],
        replacements: Optional[Any],
        kwargs: Dict[str, Any],
    ) -> List[Dict[str, Any]]:
        normalized: List[Dict[str, Any]] = []
        if isinstance(operations, list):
            normalized.extend(item for item in operations if isinstance(item, dict))
        if replacements is not None:
            normalized.append({"type": "replace_slot", "replacements": replacements})
        if kwargs.get("old_text") is not None or kwargs.get("search") is not None:
            normalized.append(
                {
                    "type": "replace_text",
                    "old_text": kwargs.get("old_text", kwargs.get("search")),
                    "new_text": kwargs.get("new_text", kwargs.get("replace", "")),
                    "slides": kwargs.get("slides"),
                }
            )
        if kwargs.get("delete_slides"):
            normalized.append({"type": "delete_slides", "slides": kwargs["delete_slides"]})
        if kwargs.get("reorder_slides"):
            normalized.append({"type": "reorder_slides", "order": kwargs["reorder_slides"]})
        return normalized

    def _resolve_path(self, path: str) -> Path:
        file_path = Path(path)
        if not file_path.is_absolute():
            file_path = self.working_dir / file_path
        return file_path.resolve()

    def _resolve_output_file(self, output_file: Optional[str], source: Path) -> Path:
        if output_file:
            path = Path(output_file)
            if not path.is_absolute():
                path = self.working_dir / path
        else:
            path = self.default_output_dir / f"{source.stem}_edited_{uuid.uuid4().hex[:8]}.pptx"
        if path.suffix.lower() != ".pptx":
            path = path.with_suffix(".pptx")
        return path.resolve()

    def get_function_schema(self) -> Dict[str, Any]:
        return {
            "name": "edit_pptx",
            "description": "编辑现有PPTX，支持replace_text、replace_slot、delete_slides、reorder_slides，并可按quality做渲染验证。",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "输入PPTX路径"},
                    "operations": {
                        "type": "array",
                        "description": "按顺序执行的操作列表。支持 {type:'replace_text', old_text, new_text, slides?}、{type:'replace_slot', slot_id, value 或 replacements}、{type:'delete_slides', slides:[2]}、{type:'reorder_slides', order:[3,1,2]}",
                        "items": {"type": "object"},
                    },
                    "replacements": {
                        "type": "object",
                        "description": "便捷参数：slot_id到新内容的映射，等价于一个 replace_slot 操作",
                    },
                    "output_file": {"type": "string", "description": "输出PPTX路径，可选"},
                    "run_validation": {"type": "boolean", "description": "是否编辑后调用validate_pptx", "default": False},
                    "quality": {
                        "type": "string",
                        "description": "质量模式：draft/standard/strict",
                        "default": "draft",
                        "enum": ["draft", "standard", "strict"],
                    },
                },
                "required": ["path"],
            },
        }

    def is_available(self) -> bool:
        return True


tool = EditPptxTool()
