"""
Production PPT generation through a PPT Master style workflow.

This tool intentionally does not use the previous DeckSpec -> PptxGenJS path.
It creates an explicit project with design/spec/page artifacts, then renders
editable PPTX shapes with python-pptx.
"""
from __future__ import annotations

import json
import re
import uuid
from collections import Counter
from html import escape
from pathlib import Path
from typing import Any, Dict, List, Optional
from urllib.parse import unquote

import structlog

from app.tools.artifact_utils import attach_document_artifact, build_artifact_resume_context
from app.tools.base.tool_interface import LLMTool, ToolCategory
from app.utils.path_config import get_data_registry, get_images_dir

logger = structlog.get_logger()


REFERENCE_DIR = Path(__file__).resolve().parent / "ppt_master_references"


def ppt_master_reference_paths() -> Dict[str, str]:
    return {
        "index": str(REFERENCE_DIR / "index.md"),
        "workflow": str(REFERENCE_DIR / "workflow.md"),
        "layout_rules": str(REFERENCE_DIR / "layout-rules.md"),
        "chart_rules": str(REFERENCE_DIR / "chart-rules.md"),
        "qa_rules": str(REFERENCE_DIR / "qa-rules.md"),
        "output_contract": str(REFERENCE_DIR / "output-contract.md"),
    }


LAYOUT_SEQUENCE = [
    "metric_strip",
    "chart_focus_callout",
    "card_grid",
    "timeline_rail",
    "matrix_summary",
    "section_quote",
]

CHART_LAYOUT_SEQUENCE = [
    "chart_focus_callout",
    "chart_left_insight_stack",
    "chart_full_bleed_insights",
    "chart_metric_sidebar",
]

STYLE_PRESETS = {
    "government_consulting": {
        "primary": "174A7C",
        "secondary": "0F766E",
        "accent": "C2410C",
        "text": "1F2937",
        "muted": "64748B",
        "surface": "F8FAFC",
        "line": "CBD5E1",
        "font": "Microsoft YaHei",
    },
    "business_clean": {
        "primary": "2563EB",
        "secondary": "0F766E",
        "accent": "DC2626",
        "text": "111827",
        "muted": "6B7280",
        "surface": "F9FAFB",
        "line": "D1D5DB",
        "font": "Microsoft YaHei",
    },
    "consulting": {
        "primary": "0F172A",
        "secondary": "334155",
        "accent": "B45309",
        "text": "111827",
        "muted": "64748B",
        "surface": "F8FAFC",
        "line": "CBD5E1",
        "font": "Microsoft YaHei",
    },
}


class CreatePptxWithPptMasterTool(LLMTool):
    def __init__(self):
        reference_paths = ppt_master_reference_paths()
        super().__init__(
            name="create_pptx_with_ppt_master",
            description=(
                "生产级PPT生成工具：按目标确认、内容结构、设计规格、母版/版式锁定、"
                "逐页绘制、质量检查、导出检查的 PPT Master 风格流程生成可编辑 PPTX。"
                "这是正式业务PPT主入口。生成 PPT 前，先读 ppt_master_references/index.md="
                f"{reference_paths['index']}，再按任务渐进读取规则。"
            ),
            category=ToolCategory.QUERY,
            version="1.0.0",
            requires_context=False,
        )
        self.default_output_dir = get_data_registry() / "presentations"
        self.default_project_root = get_data_registry() / "ppt_master_projects"

    async def execute(
        self,
        operation: str = "create",
        title: Optional[str] = None,
        purpose: str = "business_report",
        outline: Optional[List[Dict[str, Any]]] = None,
        slide_plan: Optional[List[Dict[str, Any]]] = None,
        slide_plan_path: Optional[str] = None,
        audience: str = "",
        style: str = "business_clean",
        output_file: Optional[str] = None,
        project_dir: Optional[str] = None,
        base_plan_path: Optional[str] = None,
        base_project_dir: Optional[str] = None,
        plan_patch: Optional[Dict[str, Any]] = None,
        plan_patch_path: Optional[str] = None,
        batch_slides: Optional[List[Dict[str, Any]]] = None,
        after_slide: Optional[int] = None,
        replace_slides: Optional[List[Dict[str, Any]]] = None,
        insert_slide_after: Optional[List[Dict[str, Any]]] = None,
        file_path: Optional[str] = None,
        enable_preview: bool = True,
        run_validation: bool = True,
        quality: str = "standard",
        **kwargs: Any,
    ) -> Dict[str, Any]:
        operation = str(operation or "create").strip().lower()
        if operation not in {"create", "append", "replace", "patch", "render"}:
            return {"success": False, "data": {"error": "invalid_operation"}, "summary": f"创建PPT失败：不支持 operation={operation}"}
        if operation == "render":
            return await self._render_existing_pptx(
                file_path=file_path or output_file,
                enable_preview=enable_preview,
                run_validation=run_validation,
                quality=quality,
            )

        if slide_plan_path and slide_plan is not None:
            return {"success": False, "data": {"error": "slide_plan_conflict"}, "summary": "创建PPT失败：slide_plan 和 slide_plan_path 只能传一个"}
        if plan_patch_path and plan_patch is not None:
            return {"success": False, "data": {"error": "plan_patch_conflict"}, "summary": "创建PPT失败：plan_patch 和 plan_patch_path 只能传一个"}
        if slide_plan_path:
            try:
                slide_plan = self._load_json_array_path(slide_plan_path, "slide_plan")
            except ValueError as exc:
                return {"success": False, "data": {"error": str(exc)}, "summary": f"创建PPT失败：{str(exc)[:80]}"}
        if plan_patch_path:
            try:
                plan_patch = self._load_json_object_path(plan_patch_path, "plan_patch")
            except ValueError as exc:
                return {"success": False, "data": {"error": str(exc)}, "summary": f"创建PPT失败：{str(exc)[:80]}"}

        if operation in {"append", "replace", "patch"}:
            try:
                plan_patch = self._normalize_operation_patch(
                    operation=operation,
                    plan_patch=plan_patch,
                    batch_slides=batch_slides,
                    after_slide=after_slide,
                    replace_slides=replace_slides,
                    insert_slide_after=insert_slide_after,
                )
            except ValueError as exc:
                return {"success": False, "data": {"error": str(exc)}, "summary": f"创建PPT失败：{str(exc)[:80]}"}

        is_revision = operation in {"append", "replace", "patch"} or bool(base_plan_path or base_project_dir or plan_patch)
        if not title and not is_revision:
            return {"success": False, "data": {"error": "title_required"}, "summary": "创建PPT失败：title 参数缺失"}
        if outline is not None and not isinstance(outline, list):
            return {"success": False, "data": {"error": "outline_must_be_array"}, "summary": "创建PPT失败：outline 必须是数组"}
        if slide_plan is not None and not isinstance(slide_plan, list):
            return {"success": False, "data": {"error": "slide_plan_must_be_array"}, "summary": "创建PPT失败：slide_plan 必须是数组"}
        if plan_patch is not None and not isinstance(plan_patch, dict):
            return {"success": False, "data": {"error": "plan_patch_must_be_object"}, "summary": "创建PPT失败：plan_patch 必须是对象"}

        palette = self._palette(style, kwargs.get("theme"))
        palette["font"] = self._resolve_font(palette["font"])

        revision_info: Optional[Dict[str, Any]] = None
        if is_revision:
            try:
                resolved_base_plan_path = self._resolve_base_plan_path(base_plan_path, base_project_dir)
                base_page_plan = self._load_page_plan(resolved_base_plan_path)
                page_plan = self._apply_plan_patch(base_page_plan, plan_patch or {})
                title = str(title or self._title_from_page_plan(page_plan) or "revised_presentation")
                revision_info = {
                    "base_plan_path": str(resolved_base_plan_path),
                    "patch_operation_count": self._plan_patch_operation_count(plan_patch or {}),
                }
            except ValueError as exc:
                return {"success": False, "data": {"error": str(exc)}, "summary": f"创建PPT失败：{str(exc)[:80]}"}
        else:
            title = str(title).strip()
            outline_items = self._normalize_outline(outline or [], title)
            page_plan = self._build_agent_page_plan(title, slide_plan) if slide_plan else self._build_page_plan(title, outline_items)

        output_path = self._resolve_output_file(output_file, title)
        project_path = self._resolve_project_dir(project_dir, title)
        pages_dir = project_path / "pages"
        pages_dir.mkdir(parents=True, exist_ok=True)

        design_outline = [
            {
                "title": page.get("title", f"页面 {index}"),
                "points": page.get("points", []),
                "message": page.get("message", ""),
                "role": page.get("role", "content"),
            }
            for index, page in enumerate(page_plan[1:], start=1)
        ]
        design_spec = self._build_design_spec(title, purpose, audience, style, palette, design_outline)
        self._enrich_page_plan_visuals(page_plan)
        spec_lock = self._build_spec_lock(design_spec, page_plan)

        design_spec_path = project_path / "design_spec.md"
        spec_lock_path = project_path / "spec_lock.json"
        page_plan_path = project_path / "page_plan.json"
        slide_plan_path = project_path / ("slide_plan.v2.json" if revision_info else "slide_plan.v1.json")
        design_spec_path.write_text(design_spec, encoding="utf-8")
        spec_lock_path.write_text(json.dumps(spec_lock, ensure_ascii=False, indent=2), encoding="utf-8")
        page_plan_path.write_text(json.dumps(page_plan, ensure_ascii=False, indent=2), encoding="utf-8")
        slide_plan_path.write_text(json.dumps(page_plan, ensure_ascii=False, indent=2), encoding="utf-8")

        svg_paths = []
        for page in page_plan:
            svg_path = pages_dir / f"page-{page['slide']:03d}-{page['layout']}.svg"
            svg_path.write_text(self._render_svg_preview(page, palette), encoding="utf-8")
            svg_paths.append(str(svg_path))

        self._render_pptx(output_path, title, page_plan, palette)

        result_data: Dict[str, Any] = {
            "operation": operation,
            "workflow": "ppt_master",
            "reference_paths": ppt_master_reference_paths(),
            "file_path": str(output_path),
            "output_file": str(output_path),
            "file_name": output_path.name,
            "slide_count": len(page_plan),
            "project_dir": str(project_path),
            "design_spec_path": str(design_spec_path),
            "spec_lock_path": str(spec_lock_path),
            "page_plan_path": str(page_plan_path),
            "slide_plan_path": str(slide_plan_path),
            "svg_pages": svg_paths,
            "page_plan": page_plan,
            "quality_gate": self._workflow_quality_gate(page_plan),
            "quality": quality,
        }
        if revision_info:
            result_data["revision"] = revision_info

        if enable_preview:
            try:
                from app.services.pdf_converter import pdf_converter

                result_data["pdf_preview"] = await pdf_converter.convert_to_pdf(str(output_path))
            except Exception as preview_error:
                logger.warning("ppt_master_preview_failed", error=str(preview_error))

        if run_validation or str(quality).lower() in {"standard", "strict"}:
            try:
                from app.tools.office.validate_pptx_tool import ValidatePptxTool

                validation = await ValidatePptxTool().execute(
                    str(output_path),
                    expected_fonts=[palette["font"]],
                    render_overflow_check=str(quality).lower() == "strict" or run_validation,
                )
                result_data["validation"] = validation.get("data")
                if isinstance(result_data["validation"], dict):
                    result_data["ppt_preview"] = self._ppt_preview_from_validation(result_data["validation"])
                result_data["quality_gate"] = self._workflow_quality_gate(page_plan, result_data["validation"])
            except Exception as validation_error:
                logger.warning("ppt_master_validation_failed", error=str(validation_error))
                result_data["validation_error"] = str(validation_error)
                result_data["quality_gate"] = self._workflow_quality_gate(
                    page_plan,
                    {"success": False, "issues": [{"type": "validation_error", "message": str(validation_error)}]},
                )

        final_quality_gate = result_data["quality_gate"]
        result_data["qa_status"] = final_quality_gate.get("qa_status", "passed")
        result_data["revision_tasks"] = final_quality_gate.get("revision_tasks", [])
        result_data["affected_slides"] = final_quality_gate.get("affected_slides", [])
        result_data["next_revision_base_plan_path"] = str(slide_plan_path)

        attach_document_artifact(
            result_data,
            output_path,
            kind="office",
            format="pptx",
            title=title,
            preview_key="ppt_preview" if result_data.get("ppt_preview") else "pdf_preview",
            generator=self.name,
            metadata={"workflow": "ppt_master", "project_dir": str(project_path)},
        )
        resume_context = build_artifact_resume_context(
            result_data,
            output_path,
            extra_resume={
                "project_dir": str(project_path),
                "design_spec_path": str(design_spec_path),
                "page_plan_path": str(page_plan_path),
                "slide_plan_path": str(slide_plan_path),
            },
        )
        return {
            "success": True,
            "data": result_data,
            "resources": result_data.get("resources", []),
            **resume_context,
            "summary": self._build_summary(output_path.name, len(page_plan), final_quality_gate),
        }

    async def _render_existing_pptx(
        self,
        file_path: Optional[str],
        enable_preview: bool,
        run_validation: bool,
        quality: str,
    ) -> Dict[str, Any]:
        if not file_path:
            return {"success": False, "data": {"error": "file_path_required"}, "summary": "渲染PPT失败：file_path 参数缺失"}
        pptx_path = Path(file_path)
        if not pptx_path.is_absolute():
            pptx_path = (Path.cwd() / pptx_path).resolve()
        if not pptx_path.exists():
            return {"success": False, "data": {"error": f"file_not_found: {pptx_path}"}, "summary": f"渲染PPT失败：文件不存在 {pptx_path}"}
        if pptx_path.suffix.lower() != ".pptx":
            return {"success": False, "data": {"error": "pptx_required"}, "summary": f"渲染PPT失败：只支持 .pptx，当前为 {pptx_path.suffix}"}

        result_data: Dict[str, Any] = {
            "operation": "render",
            "workflow": "ppt_master",
            "file_path": str(pptx_path),
            "output_file": str(pptx_path),
            "file_name": pptx_path.name,
            "quality": quality,
        }
        if enable_preview:
            try:
                from app.services.pdf_converter import pdf_converter

                result_data["pdf_preview"] = await pdf_converter.convert_to_pdf(str(pptx_path))
            except Exception as preview_error:
                logger.warning("ppt_master_render_preview_failed", error=str(preview_error))
        if run_validation or str(quality).lower() in {"standard", "strict"}:
            try:
                validator_cls = globals().get("ValidatePptxTool")
                if validator_cls is None:
                    from app.tools.office.validate_pptx_tool import ValidatePptxTool as validator_cls

                validation = await validator_cls().execute(
                    str(pptx_path),
                    render_overflow_check=str(quality).lower() == "strict",
                )
                result_data["validation"] = validation.get("data")
                if isinstance(result_data["validation"], dict):
                    result_data["ppt_preview"] = self._ppt_preview_from_validation(result_data["validation"])
            except Exception as validation_error:
                logger.warning("ppt_master_render_validation_failed", error=str(validation_error))
                result_data["validation_error"] = str(validation_error)

        attach_document_artifact(
            result_data,
            pptx_path,
            kind="office",
            format="pptx",
            title=pptx_path.stem,
            preview_key="ppt_preview" if result_data.get("ppt_preview") else "pdf_preview",
            generator=self.name,
            metadata={"workflow": "ppt_master", "operation": "render"},
        )
        resume_context = build_artifact_resume_context(result_data, pptx_path)
        return {
            "success": True,
            "data": result_data,
            "resources": result_data.get("resources", []),
            **resume_context,
            "summary": f"PPT渲染预览已刷新：{pptx_path.name}",
        }

    def _ppt_preview_from_validation(self, validation: Dict[str, Any]) -> Dict[str, Any]:
        return {
            "pptx_path": validation.get("pptx_path"),
            "pages": validation.get("pages", []),
            "montage_path": validation.get("montage_path"),
            "report_path": validation.get("report_path"),
        }

    def _normalize_operation_patch(
        self,
        operation: str,
        plan_patch: Optional[Dict[str, Any]],
        batch_slides: Optional[List[Dict[str, Any]]],
        after_slide: Optional[int],
        replace_slides: Optional[List[Dict[str, Any]]],
        insert_slide_after: Optional[List[Dict[str, Any]]],
    ) -> Dict[str, Any]:
        if operation == "patch" and plan_patch is not None:
            return plan_patch
        patch: Dict[str, Any] = dict(plan_patch or {})
        if operation == "append":
            if batch_slides:
                if after_slide is None:
                    raise ValueError("batch_slides_requires_after_slide")
                patch.setdefault("insert_slide_after", []).append(
                    {"after_slide": after_slide, "slides": batch_slides}
                )
            if insert_slide_after:
                patch.setdefault("insert_slide_after", []).extend(insert_slide_after)
            if not patch.get("insert_slide_after"):
                raise ValueError("append_requires_batch_slides_or_insert_slide_after")
        elif operation == "replace":
            if replace_slides:
                patch.setdefault("replace_slides", []).extend(replace_slides)
            if not patch.get("replace_slides"):
                raise ValueError("replace_requires_replace_slides")
        elif operation == "patch":
            if replace_slides:
                patch.setdefault("replace_slides", []).extend(replace_slides)
            if insert_slide_after:
                patch.setdefault("insert_slide_after", []).extend(insert_slide_after)
            if batch_slides:
                if after_slide is None:
                    raise ValueError("batch_slides_requires_after_slide")
                patch.setdefault("insert_slide_after", []).append(
                    {"after_slide": after_slide, "slides": batch_slides}
                )
            if not patch:
                raise ValueError("patch_requires_plan_patch_or_operations")
        return patch

    def _normalize_outline(self, outline: List[Dict[str, Any]], title: str) -> List[Dict[str, Any]]:
        if not outline:
            return [
                {"title": "背景与目标", "points": ["明确汇报背景", "对齐核心目标"]},
                {"title": "关键发现", "points": ["提炼主要结论", "突出证据支撑"]},
                {"title": "实施建议", "points": ["形成行动路径", "明确后续安排"]},
            ]
        normalized = []
        for index, item in enumerate(outline, start=1):
            if isinstance(item, str):
                if self._is_cover_title(item):
                    continue
                normalized.append({"title": item, "points": [], "role": self._outline_role(item)})
                continue
            if not isinstance(item, dict):
                normalized.append({"title": f"{title}要点 {index}", "points": [str(item)]})
                continue
            points = item.get("points") or item.get("bullets") or item.get("items") or []
            if isinstance(points, str):
                points = [points]
            item_title = str(item.get("title") or item.get("heading") or f"{title}要点 {index}")
            role = str(item.get("role") or self._outline_role(item_title))
            if role == "cover":
                continue
            normalized.append(
                {
                    "title": item_title,
                    "message": str(item.get("message") or item.get("summary") or ""),
                    "points": [self._point_text(point) for point in points],
                    "chart": item.get("chart"),
                    "visual": item.get("visual"),
                    "role": role,
                }
            )
        return normalized

    def _point_text(self, point: Any) -> str:
        if isinstance(point, dict):
            return str(point.get("title") or point.get("body") or point.get("text") or point)
        return str(point)

    def _build_page_plan(self, title: str, outline: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        pages = [
            {
                "slide": 1,
                "layout": "cover_statement",
                "role": "cover",
                "title": title,
                "message": "面向决策的结构化汇报",
                "points": [],
            }
        ]
        for index, item in enumerate(outline, start=2):
            layout = self._select_layout(index, item)
            pages.append(
                {
                    "slide": index,
                    "layout": layout,
                    "role": item.get("role") or "content",
                    "title": item["title"],
                    "message": item.get("message") or self._default_message(item),
                    "points": item.get("points", []),
                    "chart": item.get("chart"),
                    "visual": item.get("visual"),
                }
            )
        return pages

    def _build_agent_page_plan(self, title: str, slide_plan: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        pages = [
            {
                "slide": 1,
                "layout": "cover_statement",
                "role": "cover",
                "title": title,
                "message": "面向决策的结构化汇报",
                "points": [],
            }
        ]
        for index, item in enumerate(slide_plan, start=2):
            if not isinstance(item, dict):
                item = {"title": f"{title}页面 {index - 1}", "shapes": []}
            shapes = item.get("shapes") or item.get("shape_plan") or []
            if not isinstance(shapes, list):
                shapes = []
            pages.append(
                {
                    "slide": index,
                    "layout": "agent_shape_plan",
                    "role": item.get("role") or "content",
                    "title": str(item.get("title") or f"{title}页面 {index - 1}"),
                    "message": str(item.get("message") or ""),
                    "points": [self._point_text(point) for point in (item.get("points") or [])],
                    "chart": item.get("chart"),
                    "visual": item.get("visual"),
                    "shapes": shapes,
                }
            )
        return pages

    def _resolve_base_plan_path(self, base_plan_path: Optional[str], base_project_dir: Optional[str]) -> Path:
        if base_plan_path:
            path = Path(base_plan_path)
            if not path.is_absolute():
                path = (Path.cwd() / path).resolve()
            if not path.exists():
                raise ValueError(f"base_plan_path_not_found: {path}")
            return path
        if not base_project_dir:
            raise ValueError("base_plan_path_required")
        project = Path(base_project_dir)
        if not project.is_absolute():
            project = (Path.cwd() / project).resolve()
        if not project.exists():
            raise ValueError(f"base_project_dir_not_found: {project}")
        candidates = sorted(project.glob("slide_plan.v*.json"), key=lambda item: item.name)
        if candidates:
            return candidates[-1].resolve()
        fallback = project / "page_plan.json"
        if fallback.exists():
            return fallback.resolve()
        raise ValueError(f"base_project_plan_not_found: {project}")

    def _load_page_plan(self, path: Path) -> List[Dict[str, Any]]:
        try:
            data = json.loads(path.read_text(encoding="utf-8"))
        except json.JSONDecodeError as exc:
            raise ValueError(f"base_plan_invalid_json: {exc}") from exc
        if not isinstance(data, list):
            raise ValueError("base_plan_must_be_array")
        pages = []
        for index, item in enumerate(data, start=1):
            if not isinstance(item, dict):
                raise ValueError(f"base_plan_page_must_be_object: {index}")
            page = dict(item)
            page["slide"] = index
            pages.append(page)
        if not pages:
            raise ValueError("base_plan_empty")
        return pages

    def _resolve_json_input_path(self, path: str, field_name: str) -> Path:
        resolved = Path(path)
        if not resolved.is_absolute():
            resolved = (Path.cwd() / resolved).resolve()
        if not resolved.exists():
            raise ValueError(f"{field_name}_path_not_found: {resolved}")
        if not resolved.is_file():
            raise ValueError(f"{field_name}_path_not_file: {resolved}")
        return resolved

    def _load_json_array_path(self, path: str, field_name: str) -> List[Dict[str, Any]]:
        resolved = self._resolve_json_input_path(path, field_name)
        try:
            data = json.loads(resolved.read_text(encoding="utf-8"))
        except json.JSONDecodeError as exc:
            raise ValueError(f"{field_name}_path_invalid_json: {exc}") from exc
        if not isinstance(data, list):
            raise ValueError(f"{field_name}_path_must_be_array")
        if any(not isinstance(item, dict) for item in data):
            raise ValueError(f"{field_name}_path_items_must_be_objects")
        return data

    def _load_json_object_path(self, path: str, field_name: str) -> Dict[str, Any]:
        resolved = self._resolve_json_input_path(path, field_name)
        try:
            data = json.loads(resolved.read_text(encoding="utf-8"))
        except json.JSONDecodeError as exc:
            raise ValueError(f"{field_name}_path_invalid_json: {exc}") from exc
        if not isinstance(data, dict):
            raise ValueError(f"{field_name}_path_must_be_object")
        return data

    def _apply_plan_patch(self, base_page_plan: List[Dict[str, Any]], plan_patch: Dict[str, Any]) -> List[Dict[str, Any]]:
        pages = [dict(page) for page in base_page_plan]
        replace_slides = plan_patch.get("replace_slides") or []
        if replace_slides and not isinstance(replace_slides, list):
            raise ValueError("plan_patch.replace_slides_must_be_array")
        for operation in replace_slides:
            if not isinstance(operation, dict):
                raise ValueError("plan_patch.replace_slides_item_must_be_object")
            slide_number = self._patch_slide_number(operation)
            replacement_items = operation.get("slides")
            if replacement_items is None:
                replacement_items = operation.get("slide_plan")
            if isinstance(replacement_items, dict):
                replacement_items = [replacement_items]
            if not isinstance(replacement_items, list) or not replacement_items:
                raise ValueError("plan_patch.replace_slides_requires_slides")
            replacements = [self._normalize_patch_page(item) for item in replacement_items]
            index = self._page_index_for_slide(pages, slide_number)
            pages = pages[:index] + replacements + pages[index + 1 :]

        insertions = plan_patch.get("insert_slide_after") or []
        if isinstance(insertions, dict):
            insertions = [insertions]
        if insertions and not isinstance(insertions, list):
            raise ValueError("plan_patch.insert_slide_after_must_be_array")
        for operation in insertions:
            if not isinstance(operation, dict):
                raise ValueError("plan_patch.insert_slide_after_item_must_be_object")
            slide_number = self._patch_slide_number(operation)
            slide_items = operation.get("slides")
            if slide_items is None:
                slide_items = operation.get("slide_plan") or operation.get("page")
            if isinstance(slide_items, dict):
                slide_items = [slide_items]
            if not isinstance(slide_items, list) or not slide_items:
                raise ValueError("plan_patch.insert_slide_after_requires_slide")
            normalized_items = [self._normalize_patch_page(item) for item in slide_items]
            index = self._page_index_for_slide(pages, slide_number)
            pages = pages[: index + 1] + normalized_items + pages[index + 1 :]

        return self._renumber_page_plan(pages)

    def _patch_slide_number(self, operation: Dict[str, Any]) -> int:
        slide = operation.get("slide")
        if slide is None:
            slide = operation.get("after_slide")
        try:
            slide_number = int(slide)
        except (TypeError, ValueError) as exc:
            raise ValueError("plan_patch_slide_required") from exc
        if slide_number < 1:
            raise ValueError("plan_patch_slide_must_be_positive")
        return slide_number

    def _page_index_for_slide(self, pages: List[Dict[str, Any]], slide_number: int) -> int:
        for index, page in enumerate(pages):
            if int(page.get("slide") or index + 1) == slide_number:
                return index
        raise ValueError(f"plan_patch_slide_not_found: {slide_number}")

    def _normalize_patch_page(self, item: Any) -> Dict[str, Any]:
        if not isinstance(item, dict):
            raise ValueError("plan_patch_page_must_be_object")
        page = dict(item)
        shapes = page.get("shapes") or page.get("shape_plan") or []
        if not isinstance(shapes, list):
            shapes = []
        page["layout"] = str(page.get("layout") or "agent_shape_plan")
        page["role"] = str(page.get("role") or ("cover" if page["layout"] == "cover_statement" else "content"))
        page["title"] = str(page.get("title") or "未命名页面")
        page["message"] = str(page.get("message") or "")
        page["points"] = [self._point_text(point) for point in (page.get("points") or [])]
        page["shapes"] = shapes
        return page

    def _renumber_page_plan(self, pages: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        renumbered = []
        for index, page in enumerate(pages, start=1):
            updated = dict(page)
            updated["slide"] = index
            renumbered.append(updated)
        return renumbered

    def _title_from_page_plan(self, page_plan: List[Dict[str, Any]]) -> str:
        if page_plan:
            title = page_plan[0].get("title")
            if title:
                return str(title)
        return ""

    def _plan_patch_operation_count(self, plan_patch: Dict[str, Any]) -> int:
        total = 0
        for key in ("replace_slides", "insert_slide_after"):
            value = plan_patch.get(key)
            if isinstance(value, list):
                total += len(value)
            elif isinstance(value, dict):
                total += 1
        return total

    def _enrich_page_plan_visuals(self, page_plan: List[Dict[str, Any]]) -> None:
        for page in page_plan:
            chart = page.get("chart")
            if isinstance(chart, dict):
                asset = self._chart_asset_path(chart)
                if asset:
                    chart["resolved_asset_path"] = str(asset)
                    chart["render_mode"] = "image"
                elif self._chart_data(chart):
                    chart["render_mode"] = "native_chart"
                else:
                    chart["render_mode"] = "mock_chart"

            visual = page.get("visual")
            if isinstance(visual, dict):
                asset = self._visual_asset_path(visual)
                if asset:
                    visual["resolved_asset_path"] = str(asset)

    def _select_layout(self, index: int, item: Dict[str, Any]) -> str:
        if item.get("role") == "agenda":
            return "agenda"
        if item.get("chart"):
            return CHART_LAYOUT_SEQUENCE[(index - 2) % len(CHART_LAYOUT_SEQUENCE)]
        points = item.get("points", [])
        text = f"{item.get('title', '')} {' '.join(points)}"
        if any(keyword in text for keyword in ("流程", "步骤", "计划", "路径", "时间", "阶段")):
            return "timeline_rail"
        if any(keyword in text for keyword in ("对比", "排名", "矩阵", "差异")):
            return "matrix_summary"
        if len(points) <= 2:
            return "section_quote"
        return LAYOUT_SEQUENCE[(index - 2) % len(LAYOUT_SEQUENCE)]

    def _outline_role(self, title: str) -> str:
        text = str(title or "").strip().lower()
        if self._is_cover_title(text):
            return "cover"
        if text in {"目录", "大纲", "议程", "agenda", "contents", "table of contents"}:
            return "agenda"
        return "content"

    def _is_cover_title(self, title: str) -> bool:
        return str(title or "").strip().lower() in {"封面", "标题页", "cover", "title"}

    def _default_message(self, item: Dict[str, Any]) -> str:
        points = item.get("points") or []
        return points[0] if points else "本页围绕一个核心判断展开"

    def _build_design_spec(
        self,
        title: str,
        purpose: str,
        audience: str,
        style: str,
        palette: Dict[str, str],
        outline: List[Dict[str, Any]],
    ) -> str:
        outline_md = "\n".join(f"- {item['title']}" for item in outline)
        return (
            f"# {title} PPT Master Design Spec\n\n"
            f"## 目标\n- 用途：{purpose}\n- 受众：{audience or '未指定'}\n\n"
            f"## 风格\n- 风格：{style}\n- 主色：#{palette['primary']}\n- 强调色：#{palette['accent']}\n- 字体：{palette['font']}\n\n"
            "## 版式原则\n"
            "- 每页只讲一个核心观点。\n"
            "- 封面只出现一次，outline 中的封面/标题页必须合并到封面，不生成正文页。\n"
            "- 目录/大纲/议程必须使用 agenda 版式，不得降级为矩阵或卡片页。\n"
            "- 连续内容页不得重复同一主视觉位置。\n"
            "- 数据页必须在多种图表版式间切换，不得只在两个版式间循环。\n"
            "- 数据页优先真实图表、原生图表或大图表+洞察卡，不使用固定左图右文。\n"
            "- 流程、对比、指标、结论分别使用流程、矩阵、指标条、金句/行动页。\n\n"
            f"## 大纲\n{outline_md}\n"
        )

    def _build_spec_lock(self, design_spec: str, page_plan: List[Dict[str, Any]]) -> Dict[str, Any]:
        return {
            "workflow": "ppt_master",
            "editable_pptx": True,
            "renderer": "python-pptx-native-shapes",
            "rules": {
                "one_message_per_slide": True,
                "forbid_fixed_image_left_text_right": True,
                "min_distinct_content_layouts": 3,
            },
            "layout_sequence": [page["layout"] for page in page_plan],
            "design_spec_chars": len(design_spec),
        }

    def _workflow_quality_gate(self, page_plan: List[Dict[str, Any]], validation: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        content_layouts = [page["layout"] for page in page_plan if page["layout"] != "cover_statement"]
        repeated = sum(
            1
            for previous, current in zip(content_layouts, content_layouts[1:])
            if previous == current
        )
        distinct = len(set(content_layouts))
        issues = []
        if repeated:
            issues.append({"type": "adjacent_repeated_layout", "count": repeated})
        if distinct < min(3, len(content_layouts)):
            issues.append({"type": "low_layout_diversity", "distinct_layouts": distinct})
        if validation and validation.get("success") is False:
            issues.append({"type": "validation_failed"})
            issues.extend(self._validation_gate_issues(validation))
        issue_summary = dict(Counter(str(issue.get("type") or "unknown") for issue in issues if isinstance(issue, dict)))
        qa_failed = bool(issues) and all(
            isinstance(issue, dict) and issue.get("type") in {"validation_error", "rendered_visual_qa_unavailable"}
            for issue in issues
            if isinstance(issue, dict) and issue.get("type") != "validation_failed"
        )
        status = "qa_failed" if qa_failed else ("rewrite_required" if issues else "pass")
        qa_status = "qa_failed" if status == "qa_failed" else ("needs_revision" if issues else "passed")
        revision_tasks = [] if qa_status == "qa_failed" else self._build_revision_tasks(issues)
        affected_slides = sorted(
            {
                int(issue["slide"])
                for issue in issues
                if isinstance(issue, dict) and isinstance(issue.get("slide"), int)
            }
        )
        return {
            "status": status,
            "qa_status": qa_status,
            "rewrite_required": qa_status == "needs_revision",
            "issues": issues,
            "issue_summary": issue_summary,
            "affected_slides": affected_slides,
            "revision_tasks": revision_tasks,
            "blocking_issues": [task for task in revision_tasks if task.get("priority") == "high"],
            "layout_diversity": distinct,
        }

    def _validation_gate_issues(self, validation: Dict[str, Any]) -> List[Dict[str, Any]]:
        issues: List[Dict[str, Any]] = []
        if isinstance(validation.get("structured_issues"), list):
            issues.extend(item for item in validation["structured_issues"] if isinstance(item, dict))
            return issues
        for key in ("overflow_issues", "rendered_issues", "geometry_issues"):
            value = validation.get(key)
            if isinstance(value, list):
                issues.extend(item for item in value if isinstance(item, dict))
        for key in ("geometry", "rendered_pages", "rendered_overflow", "fonts", "design_quality", "visual_quality"):
            value = validation.get(key)
            if isinstance(value, dict) and isinstance(value.get("issues"), list):
                issues.extend(item for item in value["issues"] if isinstance(item, dict))
        if isinstance(validation.get("issues"), list):
            issues.extend(item for item in validation["issues"] if isinstance(item, dict))
        return issues

    def _build_revision_tasks(self, issues: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        tasks: List[Dict[str, Any]] = []
        seen = set()
        for issue in issues:
            if not isinstance(issue, dict):
                continue
            issue_type = str(issue.get("type") or "unknown")
            if issue_type in {"validation_failed", "validation_error", "rendered_visual_qa_unavailable"}:
                continue
            slide = issue.get("slide")
            key = (slide, issue_type)
            if key in seen:
                continue
            seen.add(key)
            task = {
                "slide": slide if isinstance(slide, int) else None,
                "type": issue_type,
                "priority": self._revision_priority(issue_type, issue),
                "message": self._revision_issue_message(issue_type, issue),
                "category": issue.get("category"),
                "location": issue.get("location") if isinstance(issue.get("location"), dict) else {},
                "evidence": issue.get("evidence") if isinstance(issue.get("evidence"), dict) else self._revision_issue_evidence(issue),
                "artifacts": issue.get("artifacts") if isinstance(issue.get("artifacts"), dict) else {},
                "issue": issue,
            }
            tasks.append(task)
        priority_order = {"high": 0, "medium": 1, "low": 2}
        return sorted(tasks, key=lambda item: (priority_order.get(item["priority"], 9), item["slide"] or 9999, item["type"]))

    def _revision_priority(self, issue_type: str, issue: Optional[Dict[str, Any]] = None) -> str:
        if isinstance(issue, dict) and issue.get("severity") in {"high", "medium", "low"}:
            return str(issue["severity"])
        if issue_type in {"shape_out_of_bounds", "rendered_content_overflow", "rendered_low_margin", "rendered_visual_overcrowding", "high_text_density"}:
            return "high"
        if issue_type in {"text_only_slide", "repeated_layout_pattern", "low_layout_diversity", "tiny_text"}:
            return "medium"
        return "low"

    def _revision_issue_message(self, issue_type: str, issue: Dict[str, Any]) -> str:
        if issue.get("message"):
            return str(issue["message"])
        messages = {
            "shape_out_of_bounds": "形状边界超出幻灯片画布。",
            "rendered_content_overflow": "渲染后页面边缘检测到非背景内容。",
            "rendered_low_margin": "渲染内容距离页面边缘过近。",
            "rendered_visual_overcrowding": "渲染视觉占用过高。",
            "high_text_density": "页面文字行数或字符数超过高密度阈值。",
            "text_only_slide": "页面主要由文本框组成，缺少图表、图片、表格或形状等视觉元素。",
            "repeated_layout_pattern": "连续页面版式重复度过高。",
            "low_layout_diversity": "整份 PPT 的内容版式多样性不足。",
            "tiny_text": "页面存在小于可读阈值的字号。",
            "expected_font_missing": "PDF 字体检测未发现期望字体。",
        }
        return messages.get(issue_type, "检测到 PPT 质量问题。")

    def _revision_issue_evidence(self, issue: Dict[str, Any]) -> Dict[str, Any]:
        excluded = {"type", "slide", "message", "category", "severity", "location", "artifacts", "raw_issue"}
        return {key: value for key, value in issue.items() if key not in excluded}

    def _revision_action(self, issue_type: str) -> str:
        actions = {
            "shape_out_of_bounds": "调整该页越界形状的位置和尺寸，确保所有元素位于画布安全边距内。",
            "rendered_content_overflow": "缩小主内容区域或拆分页面，避免渲染后内容贴边或溢出。",
            "rendered_low_margin": "增加页边距，压缩图表或底部洞察区高度，让内容离页面边缘更远。",
            "rendered_visual_overcrowding": "减少同页元素数量，优先拆页或把次要说明移入备注/下一页。",
            "high_text_density": "拆分高文字密度页面，改为主结论加卡片、图表或流程结构。",
            "text_only_slide": "补充图表、指标卡、流程、矩阵或领域示意组件，避免纯文字页。",
            "repeated_layout_pattern": "重排相邻页面版式，交替使用图表、卡片、矩阵和时间线。",
            "low_layout_diversity": "增加内容版式多样性，避免整份 PPT 只使用少数模板。",
            "tiny_text": "提高过小字号，必要时减少文本或拆页。",
            "expected_font_missing": "改用当前环境可用中文字体，或安装缺失字体后重新验证。",
        }
        return actions.get(issue_type, "根据 QA 问题调整该页版式、文字或视觉元素后重新运行 validate_pptx。")

    def _build_summary(self, output_name: str, slide_count: int, quality_gate: Dict[str, Any]) -> str:
        qa_status = quality_gate.get("qa_status", "passed")
        if qa_status == "needs_revision":
            task_count = len(quality_gate.get("revision_tasks") or [])
            return f"已生成PPT初稿：{output_name}，共 {slide_count} 页；QA发现 {task_count} 项需优化任务，建议继续迭代后交付。"
        if qa_status == "qa_failed":
            return f"已生成PPT：{output_name}，共 {slide_count} 页；QA执行异常，请先检查验证报告后再判断质量。"
        return f"已按 PPT Master 工作流生成PPT：{output_name}，共 {slide_count} 页；QA通过。"

    def _resolve_font(self, preferred: str) -> str:
        try:
            from matplotlib import font_manager

            available = {font.name for font in font_manager.fontManager.ttflist}
        except Exception:
            return preferred
        for candidate in (preferred, "Noto Sans CJK SC", "Noto Sans CJK JP", "WenQuanYi Micro Hei", "SimHei", "Arial Unicode MS"):
            if candidate in available:
                return candidate
        return preferred

    def _palette(self, style: str, theme: Optional[Dict[str, Any]]) -> Dict[str, str]:
        palette = dict(STYLE_PRESETS.get(style, STYLE_PRESETS["business_clean"]))
        for key, value in (theme or {}).items():
            mapped = {
                "primary_color": "primary",
                "secondary_color": "secondary",
                "accent_color": "accent",
                "text_color": "text",
            }.get(key, key)
            if mapped in palette and isinstance(value, str):
                palette[mapped] = value.strip("#")
        return palette

    def _render_pptx(self, output_path: Path, title: str, page_plan: List[Dict[str, Any]], palette: Dict[str, str]) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        for page in page_plan:
            slide = prs.slides.add_slide(blank)
            layout = page["layout"]
            if layout == "cover_statement":
                self._draw_cover(slide, page, palette, RGBColor, Inches, Pt, PP_ALIGN)
            elif layout == "agent_shape_plan":
                self._draw_agent_shape_plan(slide, page, palette, RGBColor, Inches, Pt)
            elif layout in set(CHART_LAYOUT_SEQUENCE):
                self._draw_chart_story(slide, page, palette, RGBColor, Inches, Pt)
            elif layout == "agenda":
                self._draw_agenda(slide, page, palette, RGBColor, Inches, Pt)
            elif layout == "metric_strip":
                self._draw_metric_strip(slide, page, palette, RGBColor, Inches, Pt)
            elif layout == "timeline_rail":
                self._draw_timeline(slide, page, palette, RGBColor, Inches, Pt)
            elif layout == "matrix_summary":
                self._draw_matrix(slide, page, palette, RGBColor, Inches, Pt)
            elif layout == "section_quote":
                self._draw_quote(slide, page, palette, RGBColor, Inches, Pt, PP_ALIGN)
            else:
                self._draw_card_grid(slide, page, palette, RGBColor, Inches, Pt)
            self._add_footer(slide, page["slide"], palette)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

    def _draw_agent_shape_plan(self, slide: Any, page: Dict[str, Any], palette: Dict[str, str], RGBColor: Any, Inches: Any, Pt: Any) -> None:
        self._add_band(slide, 0.0, 0.0, 13.333, 7.5, "FFFFFF", RGBColor, Inches)
        shapes = page.get("shapes") or []
        if not isinstance(shapes, list):
            shapes = []
        for shape in shapes:
            if not isinstance(shape, dict):
                continue
            shape_type = str(shape.get("type") or shape.get("kind") or "").lower()
            x, y, w, h = self._shape_bounds(shape)
            if w <= 0 or h <= 0:
                continue
            rendered_shape = None
            if shape_type in {"text", "textbox", "title", "body"}:
                rendered_shape = self._draw_plan_text(slide, shape, x, y, w, h, palette)
            elif shape_type in {"image", "picture"}:
                rendered_shape = self._draw_plan_image(slide, shape, x, y, w, h, RGBColor, Inches)
            elif shape_type in {"table"}:
                rendered_shape = self._draw_plan_table(slide, shape, x, y, w, h, palette, RGBColor, Inches, Pt)
            elif shape_type in {"rect", "rectangle", "card"}:
                fill = str(shape.get("fill") or shape.get("fill_color") or palette["surface"])
                line = shape.get("line") or shape.get("line_color") or palette["line"]
                rendered_shape = self._add_band(slide, x, y, w, h, fill, RGBColor, Inches, line=str(line) if line else None)
            self._apply_plan_shape_identity(rendered_shape, shape)

    def _draw_plan_text(self, slide: Any, shape: Dict[str, Any], x: float, y: float, w: float, h: float, palette: Dict[str, str]) -> Any:
        from pptx.enum.text import PP_ALIGN

        align_value = str(shape.get("align") or shape.get("alignment") or "").lower()
        align = {
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "left": PP_ALIGN.LEFT,
        }.get(align_value)
        return self._add_text(
            slide,
            str(shape.get("text") or ""),
            x,
            y,
            w,
            h,
            font_size=int(shape.get("font_size") or 14),
            color=str(shape.get("color") or palette["text"]),
            bold=bool(shape.get("bold")),
            align=align,
        )

    def _draw_plan_image(self, slide: Any, shape: Dict[str, Any], x: float, y: float, w: float, h: float, RGBColor: Any, Inches: Any) -> Any:
        path = self._resolve_asset_path(str(shape.get("path") or shape.get("image_path") or shape.get("asset") or ""))
        if not path:
            return None
        fit = str(shape.get("fit") or "contain").lower()
        if fit == "stretch":
            return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))

        image_w, image_h = self._image_size(path)
        if image_w <= 0 or image_h <= 0:
            return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        image_ratio = image_w / image_h
        box_ratio = w / h
        if fit == "cover":
            picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
            crop_left, crop_top, crop_right, crop_bottom = self._cover_crop(image_w, image_h, w, h)
            picture.crop_left = crop_left
            picture.crop_top = crop_top
            picture.crop_right = crop_right
            picture.crop_bottom = crop_bottom
            return picture
        else:
            draw_w, draw_h = (w, w / image_ratio) if image_ratio > box_ratio else (h * image_ratio, h)
        draw_x = x + (w - draw_w) / 2
        draw_y = y + (h - draw_h) / 2
        return slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), width=Inches(draw_w), height=Inches(draw_h))

    def _draw_plan_table(
        self,
        slide: Any,
        shape: Dict[str, Any],
        x: float,
        y: float,
        w: float,
        h: float,
        palette: Dict[str, str],
        RGBColor: Any,
        Inches: Any,
        Pt: Any,
    ) -> Any:
        rows = self._normalize_table_rows(shape.get("rows") or shape.get("data") or [])
        if not rows:
            return None

        row_count = len(rows)
        column_count = max(len(row) for row in rows)
        graphic_frame = slide.shapes.add_table(
            row_count,
            column_count,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        table = graphic_frame.table
        font_size = int(shape.get("font_size") or 11)
        header_fill = str(shape.get("header_fill") or shape.get("header_fill_color") or palette["primary"])
        header_color = str(shape.get("header_color") or "FFFFFF")
        cell_fill = str(shape.get("cell_fill") or shape.get("cell_fill_color") or "FFFFFF")
        text_color = str(shape.get("text_color") or shape.get("color") or palette["text"])

        for column in table.columns:
            column.width = Inches(w / column_count)
        for row in table.rows:
            row.height = Inches(h / row_count)

        for row_index, row_values in enumerate(rows):
            for column_index in range(column_count):
                cell = table.cell(row_index, column_index)
                text = row_values[column_index] if column_index < len(row_values) else ""
                fill_color = header_fill if row_index == 0 else cell_fill
                font_color = header_color if row_index == 0 else text_color
                self._format_table_cell(
                    cell,
                    text,
                    font_size=font_size,
                    color=font_color,
                    fill=fill_color,
                    bold=row_index == 0,
                    RGBColor=RGBColor,
                    Pt=Pt,
                )
        return graphic_frame

    def _normalize_table_rows(self, rows: Any) -> List[List[str]]:
        if not isinstance(rows, list):
            return []
        normalized: List[List[str]] = []
        for row in rows:
            if isinstance(row, (list, tuple)):
                normalized.append(["" if value is None else str(value) for value in row])
            elif isinstance(row, dict):
                normalized.append(["" if value is None else str(value) for value in row.values()])
            else:
                normalized.append(["" if row is None else str(row)])
        return [row for row in normalized if row]

    def _format_table_cell(
        self,
        cell: Any,
        text: str,
        *,
        font_size: int,
        color: str,
        fill: str,
        bold: bool,
        RGBColor: Any,
        Pt: Any,
    ) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = self._rgb(fill, RGBColor)
        text_frame = cell.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text
        self._set_run_font(run, "Microsoft YaHei")
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = self._rgb(color, RGBColor)

    def _set_run_font(self, run: Any, typeface: str) -> None:
        from pptx.oxml.ns import qn

        run.font.name = typeface
        run_properties = run._r.get_or_add_rPr()
        east_asian_font = run_properties.find(qn("a:ea"))
        if east_asian_font is None:
            east_asian_font = run_properties.makeelement(qn("a:ea"))
            latin_font = run_properties.find(qn("a:latin"))
            if latin_font is None:
                run_properties.append(east_asian_font)
            else:
                run_properties.insert(run_properties.index(latin_font) + 1, east_asian_font)
        east_asian_font.set("typeface", typeface)

    def _apply_plan_shape_identity(self, rendered_shape: Any, plan_shape: Dict[str, Any]) -> None:
        if rendered_shape is None:
            return
        shape_id = str(plan_shape.get("id") or plan_shape.get("shape_id") or "").strip()
        if not shape_id:
            return
        try:
            rendered_shape.name = f"pptm:{shape_id}"
        except Exception:
            logger.debug("ppt_master_shape_name_failed", shape_id=shape_id)

    def _cover_crop(self, image_w: int, image_h: int, box_w: float, box_h: float) -> tuple[float, float, float, float]:
        image_ratio = image_w / image_h
        box_ratio = box_w / box_h
        if box_ratio < image_ratio:
            crop = (1.0 - (box_ratio / image_ratio)) / 2.0
            return crop, 0.0, crop, 0.0
        if box_ratio > image_ratio:
            crop = (1.0 - (image_ratio / box_ratio)) / 2.0
            return 0.0, crop, 0.0, crop
        return 0.0, 0.0, 0.0, 0.0

    def _shape_bounds(self, shape: Dict[str, Any]) -> tuple[float, float, float, float]:
        bounds = shape.get("position") if isinstance(shape.get("position"), dict) else shape
        unit = str(bounds.get("unit") or shape.get("unit") or "in").lower()
        x = float(bounds.get("x") or bounds.get("left") or 0)
        y = float(bounds.get("y") or bounds.get("top") or 0)
        w = float(bounds.get("w") or bounds.get("width") or 0)
        h = float(bounds.get("h") or bounds.get("height") or 0)
        if unit == "relative":
            return x * 13.333, y * 7.5, w * 13.333, h * 7.5
        if unit == "emu":
            return x / 914400, y / 914400, w / 914400, h / 914400
        return x, y, w, h

    def _image_size(self, path: Path) -> tuple[int, int]:
        try:
            from PIL import Image

            with Image.open(path) as image:
                return image.size
        except Exception:
            return 0, 0

    def _add_text(self, slide: Any, text: str, x: float, y: float, w: float, h: float, *, font_size: int, color: str, bold: bool = False, align: Any = None) -> Any:
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        if align is not None:
            paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        self._set_run_font(run, "Microsoft YaHei")
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = self._rgb(color, RGBColor)
        return box

    def _draw_cover(self, slide: Any, page: Dict[str, Any], palette: Dict[str, str], RGBColor: Any, Inches: Any, Pt: Any, PP_ALIGN: Any) -> None:
        self._add_band(slide, 0.0, 0.0, 13.333, 7.5, "FFFFFF", RGBColor, Inches)
        self._add_band(slide, 0.0, 0.0, 0.28, 7.5, palette["primary"], RGBColor, Inches)
        self._add_text(slide, page["title"], 1.05, 2.35, 11.3, 0.72, font_size=30, color=palette["text"], bold=True, align=PP_ALIGN.CENTER)
        self._add_text(slide, page["message"], 1.4, 3.2, 10.6, 0.36, font_size=14, color=palette["muted"], align=PP_ALIGN.CENTER)

    def _draw_chart_story(self, slide: Any, page: Dict[str, Any], palette: Dict[str, str], RGBColor: Any, Inches: Any, Pt: Any) -> None:
        self._draw_title(slide, page, palette)
        layout = page["layout"]
        if layout == "chart_left_insight_stack":
            chart_box = (5.3, 1.55, 7.35, 4.65)
            insight_box = (0.85, 1.55, 3.65, 4.65)
        elif layout == "chart_full_bleed_insights":
            chart_box = (0.85, 1.45, 11.75, 3.75)
            insight_box = (0.85, 5.35, 11.75, 1.05)
        elif layout == "chart_metric_sidebar":
            chart_box = (3.45, 1.5, 9.05, 4.75)
            insight_box = (0.85, 1.5, 2.45, 4.75)
        else:
            chart_box = (0.85, 1.55, 7.35, 4.65)
            insight_box = (8.75, 1.55, 3.65, 4.65)
        if not self._add_real_chart_or_visual(slide, page, *chart_box, palette, RGBColor, Inches, Pt):
            self._add_mock_chart(slide, *chart_box, palette, RGBColor, Inches, Pt)
        self._add_insight_stack(slide, page, *insight_box, palette, RGBColor, Inches)

    def _add_real_chart_or_visual(
        self,
        slide: Any,
        page: Dict[str, Any],
        x: float,
        y: float,
        w: float,
        h: float,
        palette: Dict[str, str],
        RGBColor: Any,
        Inches: Any,
        Pt: Any,
    ) -> bool:
        chart = page.get("chart") if isinstance(page.get("chart"), dict) else {}
        visual = page.get("visual") if isinstance(page.get("visual"), dict) else {}

        asset = chart.get("resolved_asset_path") or visual.get("resolved_asset_path")
        if asset and Path(str(asset)).exists():
            self._add_band(slide, x, y, w, h, "FFFFFF", RGBColor, Inches, line=palette["line"])
            slide.shapes.add_picture(str(asset), Inches(x + 0.12), Inches(y + 0.12), width=Inches(w - 0.24), height=Inches(h - 0.24))
            return True

        chart_data = self._chart_data(chart)
        if chart_data:
            self._add_native_chart(slide, chart, chart_data, x, y, w, h, palette, RGBColor, Inches, Pt)
            return True

        return False

    def _add_native_chart(
        self,
        slide: Any,
        chart: Dict[str, Any],
        chart_data: Dict[str, Any],
        x: float,
        y: float,
        w: float,
        h: float,
        palette: Dict[str, str],
        RGBColor: Any,
        Inches: Any,
        Pt: Any,
    ) -> None:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

        self._add_band(slide, x, y, w, h, "FFFFFF", RGBColor, Inches, line=palette["line"])
        data = CategoryChartData()
        data.categories = chart_data["categories"]
        for series in chart_data["series"]:
            data.add_series(series["name"], series["values"])

        chart_type = str(chart.get("type") or chart.get("chart_type") or "bar").lower()
        pptx_chart_type = {
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "horizontal_bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "pie": XL_CHART_TYPE.PIE,
        }.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        if pptx_chart_type == XL_CHART_TYPE.PIE and len(chart_data["series"]) > 1:
            chart_data["series"] = chart_data["series"][:1]

        graphic_frame = slide.shapes.add_chart(pptx_chart_type, Inches(x + 0.18), Inches(y + 0.22), Inches(w - 0.36), Inches(h - 0.44), data)
        rendered_chart = graphic_frame.chart
        rendered_chart.has_legend = len(chart_data["series"]) > 1
        if rendered_chart.has_legend:
            rendered_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            rendered_chart.legend.include_in_layout = False
        if chart.get("title"):
            rendered_chart.has_title = True
            rendered_chart.chart_title.text_frame.text = str(chart["title"])

    def _draw_metric_strip(self, slide: Any, page: Dict[str, Any], palette: Dict[str, str], RGBColor: Any, Inches: Any, Pt: Any) -> None:
        self._draw_title(slide, page, palette)
        points = page.get("points") or [page["message"]]
        for index, point in enumerate(points[:5]):
            x = 0.85 + index * 2.45
            self._add_band(slide, x, 2.0, 2.05, 2.15, palette["surface"], RGBColor, Inches, line=palette["line"])
            self._add_text(slide, str(index + 1).zfill(2), x + 0.18, 2.25, 0.7, 0.35, font_size=14, color=palette["accent"], bold=True)
            self._add_text(slide, point, x + 0.18, 2.82, 1.65, 0.72, font_size=12, color=palette["text"], bold=True)
        self._add_text(slide, page["message"], 0.9, 5.25, 11.8, 0.45, font_size=18, color=palette["primary"], bold=True)

    def _draw_agenda(self, slide: Any, page: Dict[str, Any], palette: Dict[str, str], RGBColor: Any, Inches: Any, Pt: Any) -> None:
        self._draw_title(slide, page, palette)
        points = page.get("points") or ["背景", "分析", "建议"]
        for index, point in enumerate(points[:6]):
            y = 1.55 + index * 0.78
            self._add_text(slide, str(index + 1).zfill(2), 1.05, y + 0.08, 0.65, 0.28, font_size=13, color=palette["accent"], bold=True)
            self._add_band(slide, 1.9, y, 10.4, 0.48, "FFFFFF", RGBColor, Inches, line=palette["line"])
            self._add_text(slide, point, 2.15, y + 0.08, 9.6, 0.26, font_size=13, color=palette["text"], bold=True)

    def _draw_card_grid(self, slide: Any, page: Dict[str, Any], palette: Dict[str, str], RGBColor: Any, Inches: Any, Pt: Any) -> None:
        self._draw_title(slide, page, palette)
        points = page.get("points") or [page["message"]]
        for index, point in enumerate(points[:6]):
            col = index % 3
            row = index // 3
            x = 0.85 + col * 4.1
            y = 1.65 + row * 2.15
            self._add_band(slide, x, y, 3.55, 1.65, palette["surface"], RGBColor, Inches, line=palette["line"])
            self._add_text(slide, point, x + 0.25, y + 0.35, 3.05, 0.72, font_size=13, color=palette["text"], bold=True)

    def _draw_timeline(self, slide: Any, page: Dict[str, Any], palette: Dict[str, str], RGBColor: Any, Inches: Any, Pt: Any) -> None:
        self._draw_title(slide, page, palette)
        points = page.get("points") or [page["message"]]
        self._add_band(slide, 1.05, 3.25, 11.0, 0.05, palette["line"], RGBColor, Inches)
        for index, point in enumerate(points[:5]):
            x = 1.05 + index * 2.55
            self._add_band(slide, x, 2.35, 1.75, 1.75, palette["surface"], RGBColor, Inches, line=palette["line"])
            self._add_text(slide, str(index + 1), x + 0.15, 2.55, 0.35, 0.3, font_size=13, color=palette["accent"], bold=True)
            self._add_text(slide, point, x + 0.15, 3.05, 1.35, 0.62, font_size=10, color=palette["text"], bold=True)

    def _draw_matrix(self, slide: Any, page: Dict[str, Any], palette: Dict[str, str], RGBColor: Any, Inches: Any, Pt: Any) -> None:
        self._draw_title(slide, page, palette)
        points = page.get("points") or [page["message"]]
        for row in range(2):
            for col in range(2):
                index = row * 2 + col
                x = 1.05 + col * 5.7
                y = 1.75 + row * 2.1
                self._add_band(slide, x, y, 5.0, 1.55, "FFFFFF", RGBColor, Inches, line=palette["line"])
                text = points[index] if index < len(points) else page["message"]
                self._add_text(slide, text, x + 0.28, y + 0.42, 4.4, 0.55, font_size=13, color=palette["text"], bold=True)

    def _draw_quote(self, slide: Any, page: Dict[str, Any], palette: Dict[str, str], RGBColor: Any, Inches: Any, Pt: Any, PP_ALIGN: Any) -> None:
        self._draw_title(slide, page, palette)
        self._add_band(slide, 1.0, 2.15, 11.3, 2.55, palette["surface"], RGBColor, Inches)
        self._add_text(slide, page["message"], 1.55, 2.75, 10.2, 0.75, font_size=24, color=palette["primary"], bold=True, align=PP_ALIGN.CENTER)
        points = page.get("points") or []
        if len(points) > 1:
            self._add_text(slide, " / ".join(points[1:4]), 1.8, 3.75, 9.6, 0.45, font_size=12, color=palette["muted"], align=PP_ALIGN.CENTER)

    def _draw_title(self, slide: Any, page: Dict[str, Any], palette: Dict[str, str]) -> None:
        from pptx.dml.color import RGBColor
        from pptx.util import Inches

        self._add_text(slide, page["title"], 0.65, 0.42, 11.9, 0.42, font_size=22, color=palette["text"], bold=True)
        self._add_band(slide, 0.65, 1.05, 1.35, 0.06, palette["accent"], RGBColor, Inches)

    def _add_insight_stack(self, slide: Any, page: Dict[str, Any], x: float, y: float, w: float, h: float, palette: Dict[str, str], RGBColor: Any, Inches: Any) -> None:
        self._add_band(slide, x, y, w, h, palette["surface"], RGBColor, Inches, line=palette["line"])
        self._add_text(slide, page["message"], x + 0.25, y + 0.35, w - 0.5, 0.65, font_size=15, color=palette["primary"], bold=True)
        for index, point in enumerate((page.get("points") or [])[:4]):
            self._add_text(slide, point, x + 0.3, y + 1.25 + index * 0.68, w - 0.6, 0.38, font_size=10, color=palette["text"])

    def _add_mock_chart(self, slide: Any, x: float, y: float, w: float, h: float, palette: Dict[str, str], RGBColor: Any, Inches: Any, Pt: Any) -> None:
        self._add_band(slide, x, y, w, h, "FFFFFF", RGBColor, Inches, line=palette["line"])
        values = [0.45, 0.62, 0.38, 0.75, 0.58]
        for index, value in enumerate(values):
            bar_w = (w - 1.0) / len(values)
            self._add_band(slide, x + 0.45 + index * bar_w, y + h - 0.45 - value * 3.2, bar_w * 0.55, value * 3.2, palette["primary"], RGBColor, Inches)
        self._add_text(slide, "趋势 / 排名 / 结构图表区", x + 0.45, y + 0.28, w - 0.9, 0.35, font_size=11, color=palette["muted"])

    def _chart_asset_path(self, chart: Dict[str, Any]) -> Optional[Path]:
        for key in ("image_path", "path", "asset", "url", "image_url"):
            asset = chart.get(key)
            if asset:
                resolved = self._resolve_asset_path(str(asset))
                if resolved:
                    return resolved
        payload = chart.get("payload")
        if isinstance(payload, dict):
            return self._chart_asset_path(payload)
        return None

    def _visual_asset_path(self, visual: Dict[str, Any]) -> Optional[Path]:
        for key in ("asset", "image_path", "path", "url", "image_url"):
            asset = visual.get(key)
            if asset:
                resolved = self._resolve_asset_path(str(asset))
                if resolved:
                    return resolved
        return None

    def _resolve_asset_path(self, asset: str) -> Optional[Path]:
        value = str(asset or "").strip()
        if not value:
            return None

        if value.startswith("/api/image/"):
            image_id = unquote(value.removeprefix("/api/image/").split("?", 1)[0].split("#", 1)[0])
            candidate = (get_images_dir() / f"{image_id}.png").resolve()
            return candidate if candidate.exists() else None

        if value.startswith("/api/html-artifacts/"):
            parts = value.split("?", 1)[0].split("#", 1)[0].split("/")
            if len(parts) >= 5 and parts[3] and parts[4] == "assets":
                artifact_id = unquote(parts[3])
                asset_path = Path(*[unquote(part) for part in parts[5:]])
                assets_dir = (get_data_registry() / "html_artifacts" / artifact_id / "assets").resolve()
                candidate = (assets_dir / asset_path).resolve()
                try:
                    candidate.relative_to(assets_dir)
                except ValueError:
                    return None
                return candidate if candidate.exists() else None

        candidate = Path(value).expanduser()
        if not candidate.is_absolute():
            candidate = (Path.cwd() / candidate).resolve()
        return candidate if candidate.exists() else None

    def _chart_data(self, chart: Dict[str, Any]) -> Optional[Dict[str, Any]]:
        if not isinstance(chart, dict):
            return None
        payload = chart.get("payload")
        if isinstance(payload, dict):
            nested = self._chart_data(payload)
            if nested:
                return nested

        categories = chart.get("categories") or chart.get("labels") or chart.get("x")
        series = chart.get("series")
        if categories and series:
            normalized_series = self._normalize_chart_series(series)
            if normalized_series:
                return {"categories": [str(item) for item in categories], "series": normalized_series}

        rows = chart.get("data")
        if isinstance(rows, list) and rows and all(isinstance(row, dict) for row in rows):
            category_key = chart.get("category_key") or "label"
            value_key = chart.get("value_key") or "value"
            categories = [str(row.get(category_key) or row.get("name") or row.get("category") or "") for row in rows]
            values = [self._number(row.get(value_key)) for row in rows]
            if any(label for label in categories) and any(value is not None for value in values):
                return {
                    "categories": categories,
                    "series": [{"name": str(chart.get("series_name") or chart.get("title") or "数据"), "values": [value or 0 for value in values]}],
                }
        return None

    def _normalize_chart_series(self, series: Any) -> List[Dict[str, Any]]:
        if isinstance(series, dict):
            series = [series]
        if not isinstance(series, list):
            return []
        normalized = []
        for index, item in enumerate(series, start=1):
            if isinstance(item, dict):
                values = item.get("values") or item.get("data")
                name = item.get("name") or item.get("label") or f"系列{index}"
            else:
                values = item
                name = f"系列{index}"
            if not isinstance(values, list):
                continue
            number_values = [self._number(value) for value in values]
            if any(value is not None for value in number_values):
                normalized.append({"name": str(name), "values": [value or 0 for value in number_values]})
        return normalized

    def _number(self, value: Any) -> Optional[float]:
        if isinstance(value, (int, float)):
            return float(value)
        try:
            return float(str(value).replace(",", "").strip())
        except (TypeError, ValueError):
            return None

    def _add_band(self, slide: Any, x: float, y: float, w: float, h: float, fill: str, RGBColor: Any, Inches: Any, *, line: Optional[str] = None) -> Any:
        from pptx.enum.shapes import MSO_SHAPE

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(fill, RGBColor)
        if line:
            shape.line.color.rgb = self._rgb(line, RGBColor)
            shape.line.width = 1
        else:
            shape.line.fill.background()
        return shape

    def _add_footer(self, slide: Any, idx: int, palette: Dict[str, str]) -> None:
        self._add_text(slide, str(idx), 12.1, 7.08, 0.6, 0.18, font_size=8, color=palette["muted"])

    def _rgb(self, value: str, RGBColor: Any) -> Any:
        clean = str(value or "000000").strip().lstrip("#")
        if not re.fullmatch(r"[0-9a-fA-F]{6}", clean):
            clean = "000000"
        return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))

    def _render_svg_preview(self, page: Dict[str, Any], palette: Dict[str, str]) -> str:
        title = escape(page["title"])
        message = escape(page.get("message") or "")
        layout = escape(page["layout"])
        return (
            '<svg xmlns="http://www.w3.org/2000/svg" width="1600" height="900" viewBox="0 0 1600 900">'
            '<rect width="1600" height="900" fill="#ffffff"/>'
            f'<rect x="70" y="70" width="110" height="8" fill="#{palette["accent"]}"/>'
            f'<text x="70" y="145" font-size="44" font-family="Microsoft YaHei" font-weight="700" fill="#{palette["text"]}">{title}</text>'
            f'<text x="70" y="210" font-size="24" font-family="Microsoft YaHei" fill="#{palette["muted"]}">{message}</text>'
            f'<rect x="70" y="300" width="1460" height="430" fill="#{palette["surface"]}" stroke="#{palette["line"]}"/>'
            f'<text x="105" y="370" font-size="28" font-family="Microsoft YaHei" fill="#{palette["primary"]}">{layout}</text>'
            '</svg>'
        )

    def _resolve_output_file(self, output_file: Optional[str], title: str) -> Path:
        if output_file:
            path = Path(output_file)
            return path if path.is_absolute() else (Path.cwd() / path).resolve()
        safe = re.sub(r"[^\w\u4e00-\u9fff-]+", "_", title).strip("_") or "presentation"
        return self.default_output_dir / f"{safe}_{uuid.uuid4().hex[:8]}.pptx"

    def _resolve_project_dir(self, project_dir: Optional[str], title: str) -> Path:
        if project_dir:
            path = Path(project_dir)
            return path if path.is_absolute() else (Path.cwd() / path).resolve()
        safe = re.sub(r"[^\w\u4e00-\u9fff-]+", "_", title).strip("_") or "presentation"
        return self.default_project_root / f"{safe}_{uuid.uuid4().hex[:8]}"

    def get_function_schema(self) -> Dict[str, Any]:
        return {
            "name": self.name,
            "description": (
                "业务 PPT 统一入口。operation=create 新建；append/replace/patch 基于上一版 "
                "base_plan_path/base_project_dir 局部续写或修改；render 只刷新已有 PPTX 的预览/QA。"
                "生成 PPT 前必须先阅读 backend/app/tools/office/PPT操作指南.md；生成后查 qa_status/quality_gate。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "operation": {
                        "type": "string",
                        "enum": ["create", "append", "replace", "patch", "render"],
                        "default": "create",
                        "description": "create 新建；append 追加页面；replace 替换页面；patch 通用局部补丁；render 仅渲染/预览已有 PPTX。",
                    },
                    "title": {
                        "type": "string",
                        "description": "新建 PPT 标题；operation=create 时必填，续改时可省略沿用基线标题。",
                    },
                    "purpose": {
                        "type": "string",
                        "description": "用途：汇报、路演、教学、销售、总结等",
                    },
                    "audience": {"type": "string", "description": "目标受众"},
                    "style": {
                        "type": "string",
                        "enum": ["business_clean", "government_consulting", "consulting"],
                        "default": "business_clean",
                    },
                    "outline": {
                        "type": "array",
                        "description": "大纲项含 title/message/points/chart/visual；chart.image_path 生成时直接插入。",
                        "items": {"type": "object"},
                    },
                    "slide_plan": {
                        "type": "array",
                        "description": (
                            "Agent 自行规划的页面 shape plan；页含 title/message/shapes，"
                            "shape 支持 text/image/table/rect。"
                        ),
                        "items": {"type": "object"},
                    },
                    "slide_plan_path": {
                        "type": "string",
                        "description": "长 PPT 用。JSON 文件路径，内容为 slide_plan 数组。",
                    },
                    "base_plan_path": {
                        "type": "string",
                        "description": "append/replace/patch 用。上一版 data.slide_plan_path/page_plan_path。",
                    },
                    "base_project_dir": {
                        "type": "string",
                        "description": "append/replace/patch 用。上一版项目目录；自动取最新 slide_plan.v*.json。",
                    },
                    "plan_patch": {
                        "type": "object",
                        "description": "operation=patch 用。支持 replace_slides/insert_slide_after。",
                    },
                    "plan_patch_path": {
                        "type": "string",
                        "description": "长补丁用。JSON 文件路径，内容为 plan_patch 对象。",
                    },
                    "batch_slides": {
                        "type": "array",
                        "description": "operation=append/patch 的便捷追加页面数组；同时传 after_slide。",
                        "items": {"type": "object"},
                    },
                    "after_slide": {
                        "type": "integer",
                        "description": "配合 batch_slides：插入到该页之后。",
                    },
                    "replace_slides": {
                        "type": "array",
                        "description": "operation=replace/patch 的便捷替换操作数组。",
                        "items": {"type": "object"},
                    },
                    "insert_slide_after": {
                        "type": "array",
                        "description": "operation=append/patch 的插入操作数组。",
                        "items": {"type": "object"},
                    },
                    "file_path": {
                        "type": "string",
                        "description": "operation=render 用。已有 PPTX 文件路径。",
                    },
                    "output_file": {"type": "string"},
                    "quality": {"type": "string", "enum": ["draft", "standard", "strict"], "default": "standard"},
                    "run_validation": {"type": "boolean", "default": True},
                },
                "required": [],
            },
        }


tool = CreatePptxWithPptMasterTool()
